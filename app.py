import re
import os
import zipfile
import tempfile
import shutil
import magic
import random, string
import struct
from pypdf import PdfReader
from flask import Flask, render_template, request, redirect, url_for, flash, session
from sqlalchemy import text, func, extract
from sqlalchemy.orm.attributes import flag_modified
from flask_login import LoginManager, login_user, logout_user, login_required, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from werkzeug.utils import secure_filename
from datetime import datetime, timedelta, date
from itsdangerous import URLSafeTimedSerializer, BadSignature, SignatureExpired
from email_service import enviar_correo_restablecimiento, enviar_correo_confirmacion, enviar_correo_listo
from extensions import db 
from models import Usuario, Pedido, DetallePedido, ArchivoPedido, Catalogo, Configuracion, ServicioImpresion, ServicioImpresionTamano
from functools import wraps 
from pptx import Presentation
from flask_wtf.csrf import CSRFProtect
from dotenv import load_dotenv
load_dotenv()

def detectar_paginas(filepath, filename):

    ext = filename.rsplit('.', 1)[-1].lower()
    
    try:
        if ext == 'pdf':
            reader = PdfReader(filepath)
            paginas = len(reader.pages)
            return paginas, None
        
        elif ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif'):
            return 1, None
        
        elif ext == 'pptx':
            prs = Presentation(filepath)
            paginas = len(prs.slides)
            return paginas, None

        elif ext == 'zip':
            tmpdir = tempfile.mkdtemp()
            try:
                with zipfile.ZipFile(filepath, 'r') as zf:
                    # Verificar que no haya subcarpetas
                    for member in zf.namelist():
                        if member.endswith('/'):  # es una carpeta
                            return None, "El ZIP no debe contener carpetas internas."
                    zf.extractall(tmpdir)

                # Listar solo archivos en la raíz del tmpdir (sin recursión)
                archivos_internos = [f for f in os.listdir(tmpdir) if os.path.isfile(os.path.join(tmpdir, f))]
                
                imagenes_encontradas = 0
                for f in archivos_internos:
                    ext_interna = f.rsplit('.', 1)[-1].lower() if '.' in f else ''
                    if ext_interna in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif'):
                        imagenes_encontradas += 1
                    else:
                        return None, f"El ZIP contiene un archivo no soportado: {f}. Solo se permiten imágenes."

                if imagenes_encontradas == 0:
                    return None, "No se encontraron imágenes en el ZIP."
                
                return imagenes_encontradas, None
            except zipfile.BadZipFile:
                return None, "El archivo ZIP está dañado."
            except Exception as e:
                print(f"Error al procesar ZIP: {e}")
                return None, f"No se pudo leer el ZIP. ({str(e)[:80]})"
            finally:
                shutil.rmtree(tmpdir, ignore_errors=True)
        
        else:
            return None, f"Formato no soportado: .{ext}. Solo se permiten PDF, Canvas, Imágenes y ZIP de imágenes."
    
    except Exception as e:
        print(f"Error al detectar páginas de {filename}: {e}")
        return None, f"No se pudo leer el archivo. Error: {str(e)[:100]}"

def validar_mime(filepath, categorias_permitidas, filename=None):
    """
    Verifica que el archivo tenga un MIME válido según la lista de categorías.
    Si falla, intenta validar por extensión.
    Retorna (True, None) si es válido, o (False, mensaje) si no lo es.
    """
    mime = magic.from_file(filepath, mime=True)
    
    # Definir MIME permitidos (más variantes)
    mime_permitidos = {
        'pdf': ['application/pdf', 'application/x-pdf', 'application/octet-stream', 'application/vnd.pdf'],
        'imagen': ['image/png', 'image/jpeg', 'image/pjpeg', 'image/gif', 'image/bmp', 'image/tiff', 'image/x-tiff'],
        'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation', 'application/zip'],
        'zip': ['application/zip', 'application/x-zip-compressed', 'multipart/x-zip']
    }
    
    # Construir lista de MIME permitidos según categorías
    permitidos_mime = []
    for cat in categorias_permitidas:
        if cat in mime_permitidos:
            permitidos_mime.extend(mime_permitidos[cat])
    
    # Verificar por MIME
    if mime in permitidos_mime:
        return True, None
    
    # Fallback: validar por extensión (si se proporciona el nombre)
    if filename:
        ext = filename.rsplit('.', 1)[-1].lower()
        extensiones_validas = {
            'pdf': ['pdf'],
            'imagen': ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif'],
            'pptx': ['pptx'],
            'zip': ['zip']
        }
        permitidas_ext = []
        for cat in categorias_permitidas:
            permitidas_ext.extend(extensiones_validas.get(cat, []))
        if ext in permitidas_ext:
            # Log de advertencia (opcional)
            print(f"⚠️ Validación MIME falló ({mime}), pero extensión {ext} es válida. Archivo: {filename}")
            return True, None
    
    return False, f"Tipo de archivo no permitido. MIME detectado: {mime}"

def verificar_integridad(filepath, ext):
    """
    Comprueba que el archivo no está corrupto usando su librería nativa.
    Retorna (True, None) si es íntegro, o (False, mensaje) si está corrupto.
    """
    try:
        # PDF - ya validado en validar_pdf, pero intentamos abrir
        if ext == 'pdf':
            # Intentar abrir con PyPDF2 (ya lo hicimos en validar_pdf)
            reader = PdfReader(filepath)
            _ = len(reader.pages)
        
        # Imágenes
        elif ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif'):
            from PIL import Image
            with Image.open(filepath) as img:
                img.verify()  # verificar integridad
                # Volver a abrir para asegurar que se puede leer
                with Image.open(filepath) as img2:
                    img2.load()
        
        # PowerPoint
        elif ext == 'pptx':
            from pptx import Presentation
            _ = Presentation(filepath)
        
        # ZIP
        elif ext == 'zip':
            import zipfile
            with zipfile.ZipFile(filepath, 'r') as zf:
                _ = zf.namelist()
                # Verificar que el ZIP no está corrupto
                badfile = zf.testzip()
                if badfile:
                    return False, f"El ZIP contiene archivos corruptos: {badfile}"
        
        return True, None

    except Exception as e:
        return False, f"El archivo está corrupto o no se puede leer. ({str(e)[:80]})"

def validar_pdf(filepath):
    """
    Validación de PDF más permisiva: solo verifica firma y tamaño.
    Retorna (True, None) o (False, mensaje).
    """
    try:
        # Leer primeros 1024 bytes para verificar firma
        with open(filepath, 'rb') as f:
            header = f.read(1024)
        
        # 1. Verificar firma PDF (%PDF-)
        if not header.startswith(b'%PDF-'):
            return False, "El archivo no es un PDF válido (firma incorrecta)."
        
        # 2. Verificar tamaño máximo (10 MB) - opcional, ya se verifica antes
        f.seek(0, 2)
        size = f.tell()
        if size > 15 * 1024 * 1024:  # 15 MB (más permisivo)
            return False, "El archivo PDF es demasiado grande (máx. 15 MB)."
        
        # 3. Intentar leer con PyPDF2 (más fiable que buscar %%EOF)
        try:
            reader = PdfReader(filepath)
            _ = len(reader.pages)  # Si falla aquí, lanzará excepción
        except Exception as e:
            # Si falla, verificar si al menos tiene %%EOF (como último recurso)
            f.seek(max(0, size - 1024))
            trailer = f.read(1024)
            if b'%%EOF' not in trailer:
                return False, "El archivo PDF no se puede leer correctamente (posible corrupción)."
            # Si tiene EOF, pero PyPDF2 falla, permitirlo (puede ser un PDF con metadatos no estándar)
            return True, None
        
        return True, None
        
    except Exception as e:
        return False, f"No se pudo validar el PDF: {str(e)[:80]}"

def limpiar_archivos_viejos():
    """
    Elimina archivos de uploads con más de 15 días de antigüedad
    cuyo pedido esté en estado 'Entregado' o 'Cancelado'.
    Se ejecuta automáticamente al entrar al panel de admin (una vez al día).
    """
    carpetas = [
        os.path.join('static', 'uploads', 'impresion'),
        os.path.join('static', 'uploads', 'comprobantes')
    ]
    
    limite = datetime.now() - timedelta(days=15)
    eliminados = 0
    
    for carpeta in carpetas:
        ruta_completa = os.path.join(app.root_path, carpeta)
        if not os.path.exists(ruta_completa):
            continue
            
        for archivo in os.listdir(ruta_completa):
            ruta_archivo = os.path.join(ruta_completa, archivo)
            
            if os.path.isdir(ruta_archivo):
                continue
            
            # Verificar antigüedad
            fecha_mod = datetime.fromtimestamp(os.path.getmtime(ruta_archivo))
            if fecha_mod > limite:
                continue
            
            # Buscar en la base de datos
            archivo_bd = ArchivoPedido.query.filter_by(RUTA=ruta_archivo).first()
            if not archivo_bd:
                # Archivo huérfano: eliminar
                try:
                    os.remove(ruta_archivo)
                    eliminados += 1
                except:
                    pass
                continue
            
            # Verificar estado del pedido
            pedido = Pedido.query.get(archivo_bd.PEDIDO_ID)
            if pedido and pedido.ESTADO in ['Entregado', 'Cancelado']:
                try:
                    os.remove(ruta_archivo)
                    db.session.delete(archivo_bd)
                    eliminados += 1
                except:
                    pass
    
    db.session.commit()
    return eliminados

app = Flask(__name__)

UPLOAD_FOLDER = 'static/uploads'
UPLOAD_FOLDER_IMPRESION = os.path.join(UPLOAD_FOLDER, 'impresion')
UPLOAD_FOLDER_COMPROBANTES = os.path.join(UPLOAD_FOLDER, 'comprobantes')

os.makedirs(UPLOAD_FOLDER_COMPROBANTES, exist_ok=True)
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(UPLOAD_FOLDER_IMPRESION, exist_ok=True)
os.makedirs(os.path.join(UPLOAD_FOLDER_IMPRESION, 'temp'), exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024

app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get('DATABASE_URL')
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'clave-temporal-solo-dev')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
app.config['SQLALCHEMY_ENGINE_OPTIONS'] = {
    'pool_pre_ping': True,
    'pool_recycle': 300,
    'pool_timeout': 30,
}

db.init_app(app)

login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'

serializer = URLSafeTimedSerializer(app.config['SECRET_KEY'])

def admin_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or not current_user.ES_ADMIN:
            flash('Acceso denegado. Solo administradores.', 'danger')
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

def operador_required(f):
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not current_user.is_authenticated or (not current_user.ES_OPERADOR and not current_user.ES_ADMIN):
            flash('Acceso denegado. Solo operadores.', 'danger')
            return redirect(url_for('index'))
        return f(*args, **kwargs)
    return decorated_function

# ------------------------------------------------------------
# Página de inicio
# ------------------------------------------------------------

@app.route('/')
def index():
    # Obtener configuración del negocio (valores de la tabla configuracion)
    config = {}
    try:
        configuraciones = Configuracion.query.all()
        config = {c.CLAVE: c.VALOR for c in configuraciones}
    except Exception as e:
        print("Error al cargar configuración:", e)

    # Obtener imágenes del catálogo (si usas la tabla catalogo)
    try:
        with db.engine.connect() as conn:
            result = conn.execute(text("SELECT IMAGEN FROM catalogo ORDER BY ORDEN"))
            imagenes = [row[0] for row in result]
    except Exception as e:
        imagenes = []
        print("Error al consultar catálogo:", e)

    return render_template('index.html', imagenes=imagenes, config=config)

# ------------------------------------------------------------
# Login
# ------------------------------------------------------------
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.args.get('logout'):
        flash('Has cerrado sesión.', 'info')
    if request.args.get('registered'):
        flash('Cuenta creada correctamente. Ahora inicia sesión.', 'success')
    if request.args.get('reset'):
        flash('Contraseña restablecida correctamente. Inicia sesión.', 'success')
    if request.args.get('confirmed'):
        flash('Correo confirmado. Ya puedes iniciar sesión.', 'success')
    if request.method == 'POST':
        cedula = request.form.get('usuario', '').strip()
        cedula = re.sub(r'^[VEve]-', '', cedula).strip()
        contrasena = request.form.get('contrasena', '')

        error_cedula = validar_cedula(cedula)
        if error_cedula:
            flash(error_cedula, 'danger')
            return render_template('login.html')

        user = Usuario.query.filter_by(ID_USUARIO=cedula).first()

        # Verificar bloqueo
        if user and user.BLOQUEADO_HASTA and user.BLOQUEADO_HASTA > datetime.now():
            minutos_restantes = int((user.BLOQUEADO_HASTA - datetime.now()).total_seconds() // 60) + 1
            flash(f'Demasiados intentos fallidos. Espera {minutos_restantes} minuto(s) o restablece tu contraseña.', 'danger')
            return render_template('login.html')

        if user and check_password_hash(user.CONTRASEÑA, contrasena):
            # Éxito: reiniciar intentos y loguear
            user.INTENTOS_FALLIDOS = 0
            user.BLOQUEADO_HASTA = None
            db.session.commit()
            login_user(user)
            if not user.ES_ADMIN and not user.ES_OPERADOR:
                flash('¡Ahora puedes realizar peticiones en línea desde el centro de impresión!', 'info')
            flash(f'¡Bienvenido {user.NOMBRE}!', 'success')
            return redirect(url_for('index'))
        else:
            # Fallo: incrementar intentos y posible bloqueo
            if user:
                user.INTENTOS_FALLIDOS = (user.INTENTOS_FALLIDOS or 0) + 1
                if user.INTENTOS_FALLIDOS >= 5:
                    user.BLOQUEADO_HASTA = datetime.now() + timedelta(minutes=1)
                    user.INTENTOS_FALLIDOS = 0
                    db.session.commit()
                    flash('Has superado el límite de intentos. Espera 1 minuto o restablece tu contraseña.', 'danger')
                    return redirect(url_for('index', bloqueado=1))
                else:
                    db.session.commit()
            flash('Cédula o contraseña incorrectos.', 'danger')

    return render_template('login.html')
# ------------------------------------------------------------
# Registro
# ------------------------------------------------------------
@app.route('/registro', methods=['GET', 'POST'])
def registro():
    if request.method == 'POST':
        nombre = request.form.get('nombre', '').strip()
        apellido = request.form.get('apellido', '').strip()
        email = request.form.get('email', '').strip().lower()
        cedula = request.form.get('usuario_id', '').strip()
        cedula = re.sub(r'^[VEve]-?', '', cedula).strip() 
        cedula = re.sub(r'^[VEve]-', '', cedula).strip()
        telefono = request.form.get('telefono', '').strip()
        contrasena = request.form.get('contrasena', '')
        pregunta1 = request.form.get('pregunta1', '').strip()
        respuesta1 = request.form.get('respuesta1', '').strip()
        pregunta2 = request.form.get('pregunta2', '').strip()
        respuesta2 = request.form.get('respuesta2', '').strip()

        field_errors = {}

        # Validaciones existentes
        err_nombre = validar_nombre_apellido(nombre, "Nombre")
        if err_nombre:
            field_errors['nombre'] = err_nombre

        err_apellido = validar_nombre_apellido(apellido, "Apellido")
        if err_apellido:
            field_errors['apellido'] = err_apellido

        err_email = validar_email(email)
        if err_email:
            field_errors['email'] = err_email

        err_cedula = validar_cedula(cedula)
        if err_cedula:
            field_errors['usuario_id'] = err_cedula

        err_telefono = validar_telefono(telefono)
        if err_telefono:
            field_errors['telefono'] = err_telefono

        err_password = validar_contraseña(contrasena)
        if err_password:
            field_errors['contrasena'] = err_password

        # Validar preguntas de seguridad
        if not pregunta1:
            pregunta1 = '¿Nombre de tu primera mascota?'
        else:
            pregunta1 = pregunta1.strip()
            if len(pregunta1) < 4:
                field_errors['pregunta1'] = 'La pregunta debe tener al menos 4 caracteres.'
            elif len(pregunta1) > 20:
                field_errors['pregunta1'] = 'La pregunta no puede tener más de 20 caracteres.'
            elif not re.match(r'^[a-zA-ZáéíóúÁÉÍÓÚñÑ\s?]+$', pregunta1):
                field_errors['pregunta1'] = 'Solo se permiten letras y espacios.'

        if not respuesta1:
            field_errors['respuesta1'] = 'Debes escribir una respuesta.'
        else:
            respuesta1 = respuesta1.strip()
            if len(respuesta1) > 20:
                field_errors['respuesta1'] = 'La respuesta no puede tener más de 20 caracteres.'

        if not pregunta2:
            pregunta2 = '¿Ciudad donde naciste?'
        else:
            pregunta2 = pregunta2.strip()
            if len(pregunta2) < 4:
                field_errors['pregunta2'] = 'La pregunta debe tener al menos 4 caracteres.'
            elif len(pregunta2) > 20:
                field_errors['pregunta2'] = 'La pregunta no puede tener más de 20 caracteres.'
            elif not re.match(r'^[a-zA-ZáéíóúÁÉÍÓÚñÑ\s?]+$', pregunta2):
                field_errors['pregunta2'] = 'Solo se permiten letras y espacios.'

        if not respuesta2:
            field_errors['respuesta2'] = 'Debes escribir una respuesta.'
        else:
            respuesta2 = respuesta2.strip()
            if len(respuesta2) > 20:
                field_errors['respuesta2'] = 'La respuesta no puede tener más de 20 caracteres.'

        # Validar que las preguntas no sean iguales entre sí
        if pregunta1 and pregunta2 and pregunta1.lower() == pregunta2.lower():
            field_errors['pregunta2'] = 'Las preguntas no pueden ser iguales.'

        # Validar que las respuestas no sean iguales entre sí
        if respuesta1 and respuesta2 and respuesta1.lower() == respuesta2.lower():
            field_errors['respuesta2'] = 'Las respuestas no pueden ser iguales.'

        if field_errors:
            return render_template('registro.html', form_data=request.form, field_errors=field_errors)

        # Verificar duplicados (igual que antes)
        if Usuario.query.filter_by(ID_USUARIO=cedula).first():
            flash('Esta cédula ya está registrada.', 'danger')
            return render_template('registro.html', form_data=request.form, field_errors={})
        if Usuario.query.filter_by(EMAIL=email).first():
            flash('El correo ya está en uso.', 'danger')
            return render_template('registro.html', form_data=request.form, field_errors={})
        if Usuario.query.filter_by(TELEFONO=telefono).first():
            flash('El número de teléfono ya está registrado.', 'danger')
            return render_template('registro.html', form_data=request.form, field_errors={})

        nuevo = Usuario(
            ID_USUARIO=cedula,
            NOMBRE=nombre,
            APELLIDO=apellido,
            EMAIL=email,
            TELEFONO=telefono,
            CONTRASEÑA=generate_password_hash(contrasena),
            CONFIRMADO=True,   # ← confirmación directa
            PREGUNTA1=pregunta1,
            RESPUESTA1=respuesta1.strip().upper(),
            PREGUNTA2=pregunta2,
            RESPUESTA2=respuesta2.strip().upper()
        )
        db.session.add(nuevo)
        db.session.commit()
        
        return redirect(url_for('login', registered=1))

    return render_template('registro.html', form_data=None, field_errors={})

# ------------------------------------------------------------
# confirmar el correo
# ------------------------------------------------------------

@app.route('/confirmar/<token>')
def confirmar_email(token):
    try:
        user_id = serializer.loads(token, salt='confirmar-email', max_age=86400)  # 24h
    except SignatureExpired:
        flash('El enlace de confirmación ha expirado. Regístrate de nuevo.', 'danger')
        return redirect(url_for('registro'))
    except BadSignature:
        flash('Enlace inválido.', 'danger')
        return redirect(url_for('registro'))

    user = Usuario.query.get(user_id)
    if not user:
        flash('Usuario no encontrado.', 'danger')
        return redirect(url_for('registro'))

    if user.CONFIRMADO:
        flash('La cuenta ya estaba confirmada. Inicia sesión.', 'info')
    else:
        user.CONFIRMADO = True
        db.session.commit()
    return redirect(url_for('login', confirmed=1))

# ------------------------------------------------------------
# recuperar por preguntas
# ------------------------------------------------------------

@app.route('/recuperar-preguntas', methods=['GET', 'POST'])
def recuperar_preguntas():
    cedula = session.get('recuperar_cedula')
    if not cedula:
        flash('Sesión expirada. Intenta de nuevo.', 'danger')
        return redirect(url_for('recuperar'))

    user = Usuario.query.filter_by(ID_USUARIO=cedula).first()
    if not user:
        flash('Usuario no encontrado.', 'danger')
        return redirect(url_for('recuperar'))

    # Límite de intentos: 3 fallos y se bloquea
    intentos = session.get('intentos_pregunta', 0)

    if request.method == 'POST':
        respuesta = request.form.get('respuesta', '').strip().upper()
        pregunta_idx = session.get('pregunta_aleatoria')

        if pregunta_idx is None:
            flash('Error de sesión. Intenta de nuevo.', 'danger')
            return redirect(url_for('recuperar'))

        respuesta_real = user.RESPUESTA1 if pregunta_idx == 0 else user.RESPUESTA2
        if respuesta_real:
            respuesta_real = respuesta_real.strip().upper()

        if respuesta == respuesta_real:
            # Correcta: reiniciar intentos y redirigir
            session['reset_cedula'] = cedula
            session.pop('pregunta_aleatoria', None)
            session.pop('intentos_pregunta', None)
            session.pop('recuperar_cedula', None)
            return redirect(url_for('cambiar_contrasena'))
        else:
            # Fallo: alternar a la otra pregunta
            session['intentos_pregunta'] = intentos + 1
            if session['intentos_pregunta'] >= 3:
                flash('Demasiados intentos fallidos. Solicita un nuevo enlace.', 'danger')
                session.pop('pregunta_aleatoria', None)
                session.pop('intentos_pregunta', None)
                session.pop('recuperar_cedula', None)
                return redirect(url_for('recuperar'))

            # Alternar pregunta
            nuevo_idx = 1 if session['pregunta_aleatoria'] == 0 else 0
            session['pregunta_aleatoria'] = nuevo_idx
            pregunta = user.PREGUNTA1 if nuevo_idx == 0 else user.PREGUNTA2
            flash('Respuesta incorrecta. Intenta con la otra pregunta.', 'warning')
            return render_template('preguntas_seguridad.html', pregunta=pregunta)

    # GET: elegir pregunta al azar (solo la primera vez)
    if 'pregunta_aleatoria' not in session:
        import random
        session['pregunta_aleatoria'] = random.randint(0, 1)
        session['intentos_pregunta'] = 0

    idx = session['pregunta_aleatoria']
    pregunta = user.PREGUNTA1 if idx == 0 else user.PREGUNTA2

    return render_template('preguntas_seguridad.html', pregunta=pregunta)

# ------------------------------------------------------------
# Recuperar acceso
# ------------------------------------------------------------

@app.route('/recuperar', methods=['GET', 'POST'])
def recuperar():
    if request.method == 'POST':
        # --- SI YA HAY CÉDULA EN SESIÓN (usuario eligió "Enviar enlace al correo") ---
        if session.get('recuperar_cedula'):
            cedula = session['recuperar_cedula']
            user = Usuario.query.filter_by(ID_USUARIO=cedula).first()
            if user:
                try:
                    user.EMAIL.encode('ascii')
                except UnicodeEncodeError:
                    flash('El correo asociado contiene caracteres no permitidos.', 'danger')
                    return redirect(url_for('recuperar'))
                token = serializer.dumps(user.ID, salt='recuperar-contrasena')
                exito, error_msg = enviar_correo_restablecimiento(user.EMAIL, token)
                if exito:
                    flash('Te hemos enviado un enlace de recuperación a tu correo.', 'success')
                else:
                    flash(f'No se pudo enviar el correo: {error_msg}', 'danger')
                session.pop('recuperar_cedula', None)
                return redirect(url_for('login'))
            else:
                flash('Usuario no encontrado.', 'danger')
                return redirect(url_for('recuperar'))

        # --- PASO 1: Ingreso de cédula ---
        cedula = request.form.get('cedula', '').strip()
        cedula = re.sub(r'^[VEve]-', '', cedula).strip()
        error_cedula = validar_cedula(cedula)
        if error_cedula:
            flash(error_cedula, 'danger')
            return redirect(url_for('recuperar'))

        user = Usuario.query.filter_by(ID_USUARIO=cedula).first()
        if not user:
            flash('No existe una cuenta con esa cédula.', 'danger')
            return redirect(url_for('recuperar'))

        try:
            user.EMAIL.encode('ascii')
        except UnicodeEncodeError:
            flash('El correo asociado contiene caracteres no permitidos.', 'danger')
            return redirect(url_for('recuperar'))

        # Guardar cédula en sesión y mostrar opciones (SIEMPRE)
        session['recuperar_cedula'] = cedula
        tiene_preguntas = bool(user.PREGUNTA1 and user.PREGUNTA2)
        return render_template('recuperar.html', paso=2, cedula=cedula, tiene_preguntas=tiene_preguntas)

    # GET
    return render_template('recuperar.html', paso=1)

# ------------------------------------------------------------
# reset-password/<token
# ------------------------------------------------------------

@app.route('/reset-password/<token>', methods=['GET', 'POST'])
def reset_password(token):
    try:
        user_id = serializer.loads(token, salt='recuperar-contrasena', max_age=1800)
    except SignatureExpired:
        flash('El enlace ha expirado.', 'danger')
        return redirect(url_for('recuperar'))
    except BadSignature:
        flash('Enlace inválido.', 'danger')
        return redirect(url_for('recuperar'))

    user = Usuario.query.get(user_id)
    if not user:
        flash('Usuario no encontrado.', 'danger')
        return redirect(url_for('recuperar'))

    if request.method == 'POST':
        nueva = request.form.get('nueva', '')
        repetir = request.form.get('repetir', '')

        if nueva != repetir:
            flash('Las contraseñas no coinciden.', 'warning')
            return render_template('cambiar_contraseña.html', token=token)

        # Validar la nueva contraseña con los criterios definidos
        error_pass = validar_contraseña(nueva)
        if error_pass:
            flash(error_pass, 'warning')
            return render_template('cambiar_contraseña.html', token=token)

        # Actualizar la contraseña
        user.CONTRASEÑA = generate_password_hash(nueva)
        db.session.commit()
        flash('Contraseña actualizada correctamente. Inicia sesión.', 'success')

        # Limpiar sesión
        session.pop('reset_cedula', None)
        session.pop('reset_code', None)
        session.pop('reset_code_time', None)
        session.pop('recuperar_cedula', None)
        session.pop('recuperar_email_real', None)
        return redirect(url_for('login'))

    return render_template('cambiar_contraseña.html', token=token)

# ------------------------------------------------------------
# Cambiar contraseña
# ------------------------------------------------------------

@app.route('/cambiar-contrasena', methods=['GET', 'POST'])
def cambiar_contrasena():
    # Solo necesitamos la cédula en sesión
    if 'reset_cedula' not in session:
        flash('Acceso no autorizado.', 'danger')
        return redirect(url_for('recuperar'))

    cedula = session['reset_cedula']
    user = Usuario.query.filter_by(ID_USUARIO=cedula).first()
    if not user:
        flash('Usuario no encontrado.', 'danger')
        return redirect(url_for('recuperar'))

    if request.method == 'POST':
        nueva = request.form.get('nueva', '')
        repetir = request.form.get('repetir', '')

        if nueva != repetir:
            flash('Las contraseñas no coinciden.', 'warning')
        else:
            error_pass = validar_contraseña(nueva)
            if error_pass:
                flash(error_pass, 'warning')
            else:
                user.CONTRASEÑA = generate_password_hash(nueva)
                db.session.commit()

                # Limpiar TODAS las variables de sesión relacionadas con recuperación
                session.pop('reset_cedula', None)
                session.pop('reset_code', None)
                session.pop('reset_code_time', None)
                session.pop('recuperar_cedula', None)
                session.pop('recuperar_email_real', None)
                session.pop('pregunta_aleatoria', None)

                return redirect(url_for('login', reset=1))

    return render_template('cambiar_contraseña.html')

# ------------------------------------------------------------
# PANEL DE ADMINISTRACIÓN
# ------------------------------------------------------------

@app.route('/admin')
@login_required
@admin_required
def panel_admin():
    # Leer el período guardado en sesión desde Reportes/Ingresos (por defecto mensual)
    periodo = session.get('periodo_ingresos', 'mensual')
    hoy = datetime.now()

    ultima_limpieza = session.get('ultima_limpieza')
    hoy = datetime.now().date()
    
    if not ultima_limpieza or ultima_limpieza != hoy.isoformat():
        try:
            eliminados = limpiar_archivos_viejos()
            if eliminados > 0:
                print(f"[Limpieza automática] {eliminados} archivos eliminados.")
            session['ultima_limpieza'] = hoy.isoformat()
        except Exception as e:
            print(f"[Error en limpieza] {e}")
    
    # Calcular rango de fechas según el período
    if periodo == 'diario':
        desde = hoy.replace(hour=0, minute=0, second=0, microsecond=0)
        hasta = hoy.replace(hour=23, minute=59, second=59, microsecond=999999)
    elif periodo == 'semanal':
        desde = hoy - timedelta(days=hoy.weekday())
        desde = desde.replace(hour=0, minute=0, second=0, microsecond=0)
        hasta = desde + timedelta(days=6, hours=23, minutes=59, seconds=59)
    elif periodo == 'mensual':
        desde = hoy.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        if hoy.month == 12:
            hasta = hoy.replace(year=hoy.year+1, month=1, day=1) - timedelta(seconds=1)
        else:
            hasta = hoy.replace(month=hoy.month+1, day=1) - timedelta(seconds=1)
    elif periodo == 'anual':
        desde = hoy.replace(month=1, day=1, hour=0, minute=0, second=0, microsecond=0)
        hasta = hoy.replace(month=12, day=31, hour=23, minute=59, second=59, microsecond=999999)
    else:
        # Por defecto mensual
        desde = hoy.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        if hoy.month == 12:
            hasta = hoy.replace(year=hoy.year+1, month=1, day=1) - timedelta(seconds=1)
        else:
            hasta = hoy.replace(month=hoy.month+1, day=1) - timedelta(seconds=1)
    
    # Calcular ingresos SOLO de pagos confirmados y posteriores
    ingresos = db.session.query(func.sum(Pedido.TOTAL)).filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso', 'Listo', 'Entregado']),
        Pedido.FECHA >= desde,
        Pedido.FECHA <= hasta
    ).scalar() or 0
    
    pendientes = Pedido.query.filter_by(ESTADO='Pendiente de pago').count()
    nuevos_pendientes = Pedido.query.filter_by(ESTADO='Pendiente de pago', VISTO_ADMIN=False).count()
    
    en_proceso = Pedido.query.filter(
        Pedido.ESTADO.in_(['Esperando validación', 'Pago confirmado'])
    ).count()
    nuevos_proceso = Pedido.query.filter(
        Pedido.ESTADO.in_(['Esperando validación', 'Pago confirmado']),
        Pedido.VISTO_ADMIN == False
    ).count()
    
    listos = Pedido.query.filter_by(ESTADO='Listo').count()
    nuevos_listos = Pedido.query.filter_by(ESTADO='Listo', VISTO_ADMIN=False).count()  # ← NUEVO
    
    return render_template('admin/dashboard.html',
                         pendientes=pendientes,
                         proceso=en_proceso,
                         listos=listos,
                         ingresos=ingresos,
                         periodo_actual=periodo,
                         nuevos_pendientes=nuevos_pendientes,
                         nuevos_proceso=nuevos_proceso,
                         nuevos_listos=nuevos_listos)

@app.route('/admin/api/nuevos')
@login_required
@admin_required
def admin_api_nuevos():
    # Total de cada estado (para actualizar el número de la tarjeta)
    total_pendientes = Pedido.query.filter_by(ESTADO='Pendiente de pago').count()
    total_proceso = Pedido.query.filter(
        Pedido.ESTADO.in_(['Esperando validación', 'Pago confirmado'])
    ).count()
    total_listos = Pedido.query.filter_by(ESTADO='Listo').count()
    
    # Nuevos (no vistos) para los badges
    nuevos_pendientes = Pedido.query.filter_by(ESTADO='Pendiente de pago', VISTO_ADMIN=False).count()
    nuevos_proceso = Pedido.query.filter(
        Pedido.ESTADO.in_(['Esperando validación', 'Pago confirmado']),
        Pedido.VISTO_ADMIN == False
    ).count()
    nuevos_listos = Pedido.query.filter_by(ESTADO='Listo', VISTO_ADMIN=False).count()
    
    return {
        'total_pendientes': total_pendientes,
        'total_proceso': total_proceso,
        'total_listos': total_listos,
        'nuevos_pendientes': nuevos_pendientes,
        'nuevos_proceso': nuevos_proceso,
        'nuevos_listos': nuevos_listos
    }

# ------------------------------------------------------------
# PANEL ADMIN – SOLICITUDES
# ------------------------------------------------------------

@app.route('/admin/solicitudes')
@login_required
@admin_required
def admin_solicitudes():
    estado = request.args.get('estado')
    if estado:
        estados = [e.strip() for e in estado.split(',') if e.strip()]
        
        # Marcar como vistos según el estado filtrado
        if 'Pendiente de pago' in estados:
            Pedido.query.filter_by(ESTADO='Pendiente de pago', VISTO_ADMIN=False).update({'VISTO_ADMIN': True})
        if any(e in ['Esperando validación', 'Pago confirmado'] for e in estados):
            Pedido.query.filter(
                Pedido.ESTADO.in_(['Esperando validación', 'Pago confirmado']),
                Pedido.VISTO_ADMIN == False
            ).update({'VISTO_ADMIN': True})
        if 'Listo' in estados:
            Pedido.query.filter_by(ESTADO='Listo', VISTO_ADMIN=False).update({'VISTO_ADMIN': True})
        db.session.commit()
        
        if len(estados) > 1:
            pedidos = Pedido.query.filter(Pedido.ESTADO.in_(estados)).order_by(Pedido.FECHA.desc()).all()
        else:
            pedidos = Pedido.query.filter_by(ESTADO=estados[0]).order_by(Pedido.FECHA.desc()).all()
    else:
        # Sin filtro: excluir "Entregado" y "borrador"
        pedidos = Pedido.query.filter(
            ~Pedido.ESTADO.in_(['Entregado', 'borrador'])
        ).order_by(Pedido.FECHA.desc()).all()

    pedidos_con_usuarios = []
    for pedido in pedidos:
        usuario = Usuario.query.get(pedido.ID_USUARIO)
        pedidos_con_usuarios.append({'pedido': pedido, 'usuario': usuario})

    return render_template('admin/solicitudes.html',
                           pedidos=pedidos_con_usuarios,
                           estado_filtro=estado)

@app.route('/admin/solicitud/<int:pedido_id>')
@login_required
@admin_required
def admin_solicitud_detalle(pedido_id):
    pedido = db.session.get(Pedido, pedido_id) or abort(404)
    usuario = db.session.get(Usuario, pedido.ID_USUARIO)
    
    # Obtener detalles (líneas) del pedido
    detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
    archivos = ArchivoPedido.query.filter_by(PEDIDO_ID=pedido.ID).all()

    # Construir archivos_info desde los detalles (NUEVO)
    archivos_info = []
    for detalle in detalles:
        archivo = ArchivoPedido.query.get(detalle.ARCHIVO_ID) if detalle.ARCHIVO_ID else None
        item = {
            'nombre': archivo.NOMBRE_ARCHIVO if archivo else 'Sin archivo',
            'paginas': detalle.PAGINAS,
            'servicio_id': detalle.SERVICIO_ID,
            'tamano': detalle.TAMANO,
            'paginas_color': detalle.PAGINAS_COLOR,
            'comentarios': detalle.COMENTARIOS
        }
        # Enriquecer con nombre de servicio y precio
        if detalle.SERVICIO_ID:
            servicio = ServicioImpresion.query.get(detalle.SERVICIO_ID)
            item['servicio_nombre'] = servicio.TITULO if servicio else 'Desconocido'
            if detalle.TAMANO and servicio:
                tamano_obj = ServicioImpresionTamano.query.filter_by(
                    SERVICIO_ID=detalle.SERVICIO_ID, NOMBRE=detalle.TAMANO
                ).first()
                item['precio'] = float(tamano_obj.PRECIO_BN) if tamano_obj else None
        else:
            item['servicio_nombre'] = None
            item['precio'] = None
        archivos_info.append(item)

    # Marcar como visto por el admin
    if not pedido.VISTO_ADMIN:
        pedido.VISTO_ADMIN = True
        db.session.commit()

    historial = [
        {"fecha": pedido.FECHA, "estado": "Pedido creado"},
        {"fecha": datetime.now(), "estado": pedido.ESTADO}
    ]

    return render_template('admin/solicitud_detalle.html', 
                           pedido=pedido,
                           usuario=usuario,
                           detalles=detalles,
                           archivos=archivos,
                           historial=historial,
                           archivos_info=archivos_info)

@app.route('/admin/solicitud/<int:pedido_id>/confirmar-pago', methods=['POST'])
@login_required
@admin_required
def admin_confirmar_pago(pedido_id):
    pedido = db.session.get(Pedido, pedido_id)
    if pedido.ESTADO != 'Esperando validación':
        flash('El pedido no está en espera de validación.', 'warning')
        return redirect(url_for('admin_solicitud_detalle', pedido_id=pedido.ID))

    pedido.ESTADO = 'Pago confirmado'
    # Generar código de ticket si no existe
    if not pedido.CODIGO_TICKET:
        codigo = f"TICK-{pedido.ID}-{''.join(random.choices(string.ascii_uppercase + string.digits, k=4))}"
        pedido.CODIGO_TICKET = codigo
    db.session.commit()
    flash('Pago confirmado. El cliente ya puede acceder a su ticket.', 'success')
    return redirect(url_for('admin_solicitud_detalle', pedido_id=pedido.ID))

@app.route('/admin/solicitud/<int:pedido_id>/actualizar', methods=['POST'])
@login_required
@admin_required
def admin_actualizar_estado(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    nuevo_estado = request.form.get('estado')
    
    # Validar que el nuevo estado sea permitido
    if nuevo_estado not in ['Listo', 'Entregado']:
        flash('Estado no válido. Solo se permiten "Listo" o "Entregado".', 'danger')
        return redirect(url_for('admin_solicitud_detalle', pedido_id=pedido.ID))
    
    # Validar transición permitida según estado actual
    if nuevo_estado == 'Listo' and pedido.ESTADO not in ['Pago confirmado', 'En proceso']:
        flash('Solo se puede marcar como "Listo" si el pedido está en "Pago confirmado" o "En proceso".', 'danger')
        return redirect(url_for('admin_solicitud_detalle', pedido_id=pedido.ID))
    
    if nuevo_estado == 'Entregado' and pedido.ESTADO != 'Listo':
        flash('Solo se puede marcar como "Entregado" si el pedido está en "Listo".', 'danger')
        return redirect(url_for('admin_solicitud_detalle', pedido_id=pedido.ID))
    
    # Actualizar estado
    pedido.ESTADO = nuevo_estado
    db.session.commit()
    
    # Si el nuevo estado es Listo, enviar correo al cliente
    if nuevo_estado == 'Listo':
        usuario = Usuario.query.get(pedido.ID_USUARIO)
        if usuario:
            ticket_url = url_for('ticket_impresion', pedido_id=pedido.ID, _external=True)
            if not pedido.CODIGO_TICKET:
                import random, string
                codigo = f"TICK-{pedido.ID}-{''.join(random.choices(string.ascii_uppercase + string.digits, k=4))}"
                pedido.CODIGO_TICKET = codigo
                db.session.commit()
            exito, error = enviar_correo_listo(usuario.EMAIL, usuario.NOMBRE, ticket_url)
            if exito:
                flash('Estado actualizado y correo enviado al cliente.', 'success')
            else:
                flash(f'Estado actualizado pero no se pudo enviar el correo: {error}', 'warning')
        else:
            flash('Estado actualizado.', 'success')
    else:
        flash('Estado actualizado correctamente.', 'success')
    
    return redirect(url_for('admin_solicitud_detalle', pedido_id=pedido.ID))

@app.route('/admin/solicitud/<int:pedido_id>/eliminar', methods=['POST'])
@login_required
@admin_required
def admin_eliminar_solicitud(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    
    # Eliminar primero los detalles asociados
    DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).delete()
    
    # Eliminar archivos asociados
    ArchivoPedido.query.filter_by(PEDIDO_ID=pedido.ID).delete()
    
    # Eliminar el pedido
    db.session.delete(pedido)
    db.session.commit()
    
    flash('Solicitud eliminada correctamente.', 'success')
    return redirect(url_for('admin_solicitudes'))

# ------------------------------------------------------------
# PANEL ADMIN – CATALOGO
# ------------------------------------------------------------

@app.route('/admin/catalogos')
@login_required
@admin_required
def admin_catalogos():
    tab = request.args.get('tab', 'servicios')
    imagenes = Catalogo.query.order_by(Catalogo.ORDEN).all()          
    servicios = ServicioImpresion.query.order_by(ServicioImpresion.TITULO).all()

    return render_template('admin/catalogos.html', 
                           imagenes=imagenes,
                           servicios=servicios,
                           tab_activa=tab)

@app.route('/admin/catalogos/agregar', methods=['POST'])
@login_required
@admin_required
def admin_agregar_imagen():
    imagen = request.files.get('imagen')
    if not imagen or imagen.filename == '':
        flash('Debes seleccionar una imagen.', 'danger')
        return redirect(url_for('admin_catalogos'))

    filename = secure_filename(imagen.filename)
    filepath = os.path.join('static/img', filename)
    imagen.save(filepath)

    # Verificar si ya existe en la base de datos
    if Catalogo.query.filter_by(IMAGEN=filename).first():
        flash('Ya existe una imagen con ese nombre.', 'warning')
        return redirect(url_for('admin_catalogos'))

    nuevo = Catalogo(IMAGEN=filename)
    db.session.add(nuevo)
    db.session.commit()
    flash('Imagen agregada correctamente.', 'success')
    return redirect(url_for('admin_catalogos'))

@app.route('/admin/catalogos/eliminar/<int:imagen_id>', methods=['POST'])
@login_required
@admin_required
def admin_eliminar_imagen(imagen_id):
    imagen = Catalogo.query.get_or_404(imagen_id)
    # Eliminar archivo físico
    ruta = os.path.join('static/img', imagen.IMAGEN)
    if os.path.exists(ruta):
        os.remove(ruta)
    db.session.delete(imagen)
    db.session.commit()
    flash('Imagen eliminada correctamente.', 'success')
    return redirect(url_for('admin_catalogos'))

@app.route('/admin/catalogos/actualizar/<int:imagen_id>', methods=['POST'])
@login_required
@admin_required
def admin_actualizar_imagen(imagen_id):
    imagen = Catalogo.query.get_or_404(imagen_id)
    nueva_imagen = request.files.get('imagen')
    if nueva_imagen and nueva_imagen.filename:
        # Eliminar la anterior
        ruta_anterior = os.path.join('static/img', imagen.IMAGEN)
        if os.path.exists(ruta_anterior):
            os.remove(ruta_anterior)
        # Guardar la nueva
        filename = secure_filename(nueva_imagen.filename)
        nueva_imagen.save(os.path.join('static/img', filename))
        imagen.IMAGEN = filename
        db.session.commit()
        flash('Imagen actualizada correctamente.', 'success')
    else:
        flash('No se seleccionó una nueva imagen.', 'info')
    return redirect(url_for('admin_catalogos'))

@app.route('/admin/catalogos/orden-actual')
@login_required
@admin_required
def admin_orden_actual():
    imagenes = Catalogo.query.order_by(Catalogo.ORDEN).all()
    orden = [{'id': img.ID, 'nombre': img.IMAGEN} for img in imagenes]
    return {'orden': orden}

@app.route('/admin/catalogos/reordenar', methods=['POST'])
@login_required
@admin_required
def admin_reordenar_catalogo():
    data = request.get_json()
    orden_ids = data.get('orden', [])
    for idx, img_id in enumerate(orden_ids):
        Catalogo.query.filter_by(ID=img_id).update({'ORDEN': idx})
    db.session.commit()
    return {'success': True}

# SERVICIOS DE IMPRESIÓN
@app.route('/admin/servicios-impresion')
@login_required
@admin_required
def admin_servicios_impresion():
    servicios = ServicioImpresion.query.order_by(ServicioImpresion.TITULO).all()
    return render_template('admin/servicios_impresion.html', servicios=servicios)

@app.route('/admin/servicio-impresion/agregar', methods=['POST'])
@login_required
@admin_required
def admin_agregar_servicio_impresion():
    titulo = request.form.get('titulo', '').strip()
    descripcion = request.form.get('descripcion', '').strip()
    activo = 'activo' in request.form
    es_mixto = 'es_mixto' in request.form

    if not titulo:
        flash('El título es obligatorio.', 'danger')
        return redirect(url_for('admin_catalogos', tab='serviciosImpresion'))

    nuevo = ServicioImpresion(TITULO=titulo, 
                              DESCRIPCION=descripcion, 
                              ACTIVO=activo,
                              ES_MIXTO=es_mixto)
    
    db.session.add(nuevo)
    db.session.flush() 

    # Procesar tamaños
    nombres = request.form.getlist('tam_nombre[]')
    precios_bn = request.form.getlist('tam_precio_bn[]')
    precios_color = request.form.getlist('tam_precio_color[]')
    
    # Iterar con índice para poder acceder al precio color correcto
    for i, nom in enumerate(nombres):
        pre_bn = precios_bn[i] if i < len(precios_bn) else ''
        if nom.strip() and pre_bn.strip():
            try:
                precio_bn = float(pre_bn)
            except ValueError:
                flash('Precio B/N inválido.', 'danger')
                db.session.rollback()
                return redirect(url_for('admin_catalogos', tab='serviciosImpresion'))
            
            # Precio color: tomar el correspondiente, o si no existe, usar el B/N
            pre_color = precios_color[i] if i < len(precios_color) else pre_bn
            try:
                precio_color = float(pre_color)
            except ValueError:
                precio_color = precio_bn
            
            t = ServicioImpresionTamano(
                SERVICIO_ID=nuevo.ID,
                NOMBRE=nom.strip(),
                PRECIO_BN=precio_bn,
                PRECIO_COLOR=precio_color
            )
            db.session.add(t)

    db.session.commit()
    flash('Servicio de impresión agregado correctamente.', 'success')
    return redirect(url_for('admin_catalogos', tab='serviciosImpresion'))

@app.route('/admin/servicio-impresion/<int:servicio_id>/editar', methods=['POST'])
@login_required
@admin_required
def admin_editar_servicio_impresion(servicio_id):
    servicio = ServicioImpresion.query.get_or_404(servicio_id)
    servicio.TITULO = request.form.get('titulo', '').strip()
    servicio.DESCRIPCION = request.form.get('descripcion', '').strip()
    servicio.ACTIVO = 'activo' in request.form
    servicio.ES_MIXTO = 'es_mixto' in request.form

    # Reemplazar tamaños existentes por los nuevos enviados
    ServicioImpresionTamano.query.filter_by(SERVICIO_ID=servicio_id).delete()
    nombres = request.form.getlist('tam_nombre[]')
    precios_bn = request.form.getlist('tam_precio_bn[]')
    precios_color = request.form.getlist('tam_precio_color[]')
    
    for i, nom in enumerate(nombres):
        pre_bn = precios_bn[i] if i < len(precios_bn) else ''
        if nom.strip() and pre_bn.strip():
            pre_color = precios_color[i] if i < len(precios_color) else pre_bn
            try:
                precio_color = float(pre_color)
            except ValueError:
                precio_color = float(pre_bn)
            t = ServicioImpresionTamano(
                SERVICIO_ID=servicio_id,
                NOMBRE=nom.strip(),
                PRECIO_BN=float(pre_bn),
                PRECIO_COLOR=precio_color
            )
            db.session.add(t)
            
    db.session.commit()

    flash('Servicio actualizado.', 'success')
    return redirect(url_for('admin_catalogos', tab='serviciosImpresion'))

@app.route('/admin/servicio-impresion/<int:servicio_id>/eliminar', methods=['POST'])
@login_required
@admin_required
def admin_eliminar_servicio_impresion(servicio_id):
    servicio = ServicioImpresion.query.get_or_404(servicio_id)
    
    # Desvincular pedidos que usan este servicio (poner SERVICIO_ID a NULL)
    Pedido.query.filter_by(SERVICIO_ID=servicio_id).update({'SERVICIO_ID': None})
    
    db.session.delete(servicio)
    db.session.commit()
    flash('Servicio de impresión eliminado.', 'success')
    return redirect(url_for('admin_catalogos', tab='serviciosImpresion'))

# ------------------------------------------------------------
# PANEL ADMIN – USUARIOS
# ------------------------------------------------------------

@app.route('/admin/usuarios')
@login_required
@admin_required
def admin_usuarios():
    busqueda = request.args.get('buscar', '').strip()
    if busqueda:
        # Buscar por cédula, nombre, apellido o email
        usuarios = Usuario.query.filter(
            (Usuario.ID_USUARIO.contains(busqueda)) |
            (Usuario.NOMBRE.contains(busqueda)) |
            (Usuario.APELLIDO.contains(busqueda)) |
            (Usuario.EMAIL.contains(busqueda))
        ).order_by(Usuario.APELLIDO, Usuario.NOMBRE).all()
    else:
        usuarios = Usuario.query.order_by(Usuario.APELLIDO, Usuario.NOMBRE).all()
    return render_template('admin/usuarios.html', usuarios=usuarios, busqueda=busqueda)

@app.route('/admin/usuario/<int:user_id>')
@login_required
@admin_required
def admin_usuario_detalle(user_id):
    usuario = Usuario.query.get_or_404(user_id)
    pedidos = Pedido.query.filter_by(ID_USUARIO=usuario.ID).order_by(Pedido.FECHA.desc()).all()
    return render_template('admin/usuario_detalle.html', usuario=usuario, pedidos=pedidos)

@app.route('/admin/usuario/<int:user_id>/eliminar', methods=['POST'])
@login_required
@admin_required
def admin_eliminar_usuario(user_id):
    usuario = Usuario.query.get_or_404(user_id)
    
    if usuario.ID == current_user.ID:
        flash('No puedes eliminar tu propio usuario.', 'danger')
        return redirect(url_for('admin_usuarios'))
    
    # Eliminar pedidos asociados (y sus detalles, archivos)
    pedidos = Pedido.query.filter_by(ID_USUARIO=usuario.ID).all()
    for pedido in pedidos:
        DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).delete()
        ArchivoPedido.query.filter_by(PEDIDO_ID=pedido.ID).delete()
        db.session.delete(pedido)
    
    # Ya no se eliminan mensajes de chat porque la tabla fue borrada
    
    db.session.delete(usuario)
    db.session.commit()
    
    flash('Usuario eliminado correctamente.', 'success')
    return redirect(url_for('admin_usuarios'))

# ------------------------------------------------------------
# ADMIN - OPERADORES
# ------------------------------------------------------------
@app.route('/admin/operadores')
@login_required
@admin_required
def admin_operadores():
    # Obtener todos los usuarios que ya son operadores
    operadores = Usuario.query.filter_by(ES_OPERADOR=1).all()
    return render_template('admin/operadores.html', operadores=operadores)

@app.route('/admin/operadores/buscar', methods=['POST'])
@login_required
@admin_required
def admin_operadores_buscar():
    cedula = request.form.get('cedula', '').strip()
    cedula = re.sub(r'^[VEve]-', '', cedula).strip()
    if not cedula:
        return {'error': 'Debe ingresar una cédula'}, 400
    
    usuario = Usuario.query.filter_by(ID_USUARIO=cedula).first()
    if not usuario:
        return {'error': 'No se encontró un usuario con esa cédula'}, 404

    # Contar cuántos operadores hay actualmente
    cantidad_operadores = Usuario.query.filter_by(ES_OPERADOR=1).count()
    
    return {
        'id': usuario.ID,
        'nombre': usuario.NOMBRE,
        'apellido': usuario.APELLIDO,
        'cedula': usuario.ID_USUARIO,
        'email': usuario.EMAIL,
        'telefono': usuario.TELEFONO,
        'es_operador': usuario.ES_OPERADOR,
        'cantidad_operadores': cantidad_operadores
    }

@app.route('/admin/operadores/asignar', methods=['POST'])
@login_required
@admin_required
def admin_operadores_asignar():
    user_id = request.form.get('user_id')
    if not user_id:
        flash('ID de usuario no proporcionado', 'danger')
        return redirect(url_for('admin_operadores'))
    
    usuario = Usuario.query.get(int(user_id))
    if not usuario:
        flash('Usuario no encontrado', 'danger')
        return redirect(url_for('admin_operadores'))
    
    # Si ya es operador, lo desasignamos
    if usuario.ES_OPERADOR:
        usuario.ES_OPERADOR = False
        db.session.commit()
        flash(f'{usuario.NOMBRE} {usuario.APELLIDO} ya no es operador.', 'success')
    else:
        # Verificar límite de 6 operadores
        if Usuario.query.filter_by(ES_OPERADOR=1).count() >= 6:
            flash('No se pueden agregar más de 6 operadores.', 'danger')
        else:
            usuario.ES_OPERADOR = True
            db.session.commit()
            flash(f'{usuario.NOMBRE} {usuario.APELLIDO} ahora es operador.', 'success')
    
    return redirect(url_for('admin_operadores'))

# ------------------------------------------------------------
# PANEL ADMIN – REPORTES
# ------------------------------------------------------------

@app.route('/admin/reportes')
@login_required
@admin_required
def admin_reportes():
    # --- Parámetros de período y pestaña (ya existentes) ---
    periodo = request.args.get('periodo', session.get('periodo_ingresos', 'mensual'))
    session['periodo_ingresos'] = periodo
    tab_activa = request.args.get('tab', 'ventas')
    desde_fecha = request.args.get('desde', '')
    hasta_fecha = request.args.get('hasta', '')

    query = Pedido.query

    if desde_fecha:
        query = query.filter(Pedido.FECHA >= datetime.strptime(desde_fecha, '%Y-%m-%d'))
    if hasta_fecha:
        query = query.filter(Pedido.FECHA <= datetime.strptime(hasta_fecha, '%Y-%m-%d') + timedelta(days=1))

    # EXCLUIR BORRADORES en solicitudes por fecha
    solicitudes_fecha = query.filter(~Pedido.ESTADO.in_(['borrador'])).order_by(Pedido.FECHA.desc()).all()
    buscar = request.args.get('buscar', '').strip()

    hoy = datetime.now()

    # Calcular rango de fechas según período (para ingresos)
    if periodo == 'diario':
        desde = hoy.replace(hour=0, minute=0, second=0, microsecond=0)
        hasta = hoy.replace(hour=23, minute=59, second=59, microsecond=999999)
    elif periodo == 'semanal':
        desde = hoy - timedelta(days=hoy.weekday())
        desde = desde.replace(hour=0, minute=0, second=0, microsecond=0)
        hasta = desde + timedelta(days=6, hours=23, minutes=59, seconds=59)
    elif periodo == 'mensual':
        desde = hoy.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        if hoy.month == 12:
            hasta = hoy.replace(year=hoy.year+1, month=1, day=1) - timedelta(seconds=1)
        else:
            hasta = hoy.replace(month=hoy.month+1, day=1) - timedelta(seconds=1)
    elif periodo == 'anual':
        desde = hoy.replace(month=1, day=1, hour=0, minute=0, second=0, microsecond=0)
        hasta = hoy.replace(month=12, day=31, hour=23, minute=59, second=59, microsecond=999999)
    else:
        desde = hoy.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
        if hoy.month == 12:
            hasta = hoy.replace(year=hoy.year+1, month=1, day=1) - timedelta(seconds=1)
        else:
            hasta = hoy.replace(month=hoy.month+1, day=1) - timedelta(seconds=1)

    # --- Servicios más solicitados (ranking) ---
    # EXCLUIR BORRADORES en el ranking
    servicios_ranking = db.session.query(
        ServicioImpresion.ID,
        ServicioImpresion.TITULO,
        ServicioImpresion.DESCRIPCION,
        func.count(DetallePedido.ID).label('total_pedidos')
    ).join(DetallePedido, DetallePedido.SERVICIO_ID == ServicioImpresion.ID)\
     .join(Pedido, Pedido.ID == DetallePedido.PEDIDO_ID)\
     .filter(Pedido.ESTADO.in_(['Pago confirmado', 'En proceso', 'Listo', 'Entregado']))\
     .filter(Pedido.ESTADO != 'borrador')\
     .group_by(ServicioImpresion.ID)\
     .order_by(func.count(DetallePedido.ID).desc())\
     .all()

    ranking = []
    for idx, (id_srv, titulo, descripcion, total) in enumerate(servicios_ranking, start=1):
        ranking.append({
            'posicion': idx,
            'titulo': titulo,
            'descripcion': descripcion,
            'total': total
        })

    # Ingresos del período
    ingresos_periodo = db.session.query(func.sum(Pedido.TOTAL)).filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso', 'Listo', 'Entregado']),
        Pedido.ESTADO != 'borrador',
        Pedido.FECHA >= desde,
        Pedido.FECHA <= hasta
    ).scalar() or 0

    # --- Ventas con búsqueda y datos de usuario ---
    # Consulta base
    ventas_query = Pedido.query.filter(Pedido.ESTADO != 'borrador')  # EXCLUIR BORRADORES

    # Si hay búsqueda, filtrar por datos del usuario
    if buscar:
        usuarios_ids = Usuario.query.filter(
            (Usuario.ID_USUARIO.contains(buscar)) |
            (Usuario.NOMBRE.contains(buscar)) |
            (Usuario.APELLIDO.contains(buscar)) |
            (Usuario.EMAIL.contains(buscar))
        ).with_entities(Usuario.ID).all()
        ids = [u[0] for u in usuarios_ids]
        if ids:
            ventas_query = ventas_query.filter(Pedido.ID_USUARIO.in_(ids))
        else:
            ventas_query = ventas_query.filter(False)  # sin resultados

    ventas = ventas_query.order_by(Pedido.FECHA.desc()).all()

    # Añadir usuario a cada pedido
    ventas_con_usuarios = []
    for pedido in ventas:
        usuario = Usuario.query.get(pedido.ID_USUARIO)
        ventas_con_usuarios.append({'pedido': pedido, 'usuario': usuario})

    # Solicitudes por fecha (con filtro de fechas) - ya excluye borradores
    query_fechas = Pedido.query.filter(Pedido.ESTADO != 'borrador')
    if desde_fecha:
        query_fechas = query_fechas.filter(Pedido.FECHA >= datetime.strptime(desde_fecha, '%Y-%m-%d'))
    if hasta_fecha:
        query_fechas = query_fechas.filter(Pedido.FECHA <= datetime.strptime(hasta_fecha, '%Y-%m-%d') + timedelta(days=1))
    solicitudes_fecha = query_fechas.order_by(Pedido.FECHA.desc()).all()

    # Ingresos mensuales (histórico)
    ingresos_mensuales = db.session.query(
        extract('year', Pedido.FECHA).label('anio'),
        extract('month', Pedido.FECHA).label('mes'),
        func.sum(Pedido.TOTAL).label('total')
    ).filter(Pedido.ESTADO.in_(['Pago confirmado', 'En proceso', 'Listo', 'Entregado']))\
     .filter(Pedido.ESTADO != 'borrador')\
     .group_by('anio', 'mes')\
     .order_by('anio', 'mes').all()

    filtro_activo = bool(desde_fecha or hasta_fecha)

    return render_template('admin/reportes.html',
                           ventas=ventas_con_usuarios,
                           solicitudes_fecha=solicitudes_fecha,
                           ingresos_mensuales=ingresos_mensuales,
                           ingresos_periodo=ingresos_periodo,
                           periodo_actual=periodo,
                           desde=desde_fecha,
                           hasta=hasta_fecha,
                           tab_activa=tab_activa,
                           buscar=buscar,
                           filtro_activo=filtro_activo,
                           hoy=datetime.now(),
                           ranking=ranking)

@app.route('/admin/api/pedido/<int:pedido_id>')
@login_required
@admin_required
def admin_api_pedido(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    usuario = Usuario.query.get(pedido.ID_USUARIO)
    
    # Obtener detalles (líneas)
    detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido_id).all()
    archivos_info = []
    total_paginas = 0
    detalles_facturacion = []

    if detalles:
        for detalle in detalles:
            archivo = ArchivoPedido.query.get(detalle.ARCHIVO_ID) if detalle.ARCHIVO_ID else None
            
            # Obtener servicio y precio
            servicio_nombre = None
            precio = None
            if detalle.SERVICIO_ID:
                servicio = ServicioImpresion.query.get(detalle.SERVICIO_ID)
                if servicio:
                    servicio_nombre = servicio.TITULO
                    if detalle.TAMANO:
                        tamano_obj = ServicioImpresionTamano.query.filter_by(
                            SERVICIO_ID=detalle.SERVICIO_ID, NOMBRE=detalle.TAMANO
                        ).first()
                        if tamano_obj:
                            precio = float(tamano_obj.PRECIO_BN)
            
            archivos_info.append({
                'nombre': archivo.NOMBRE_ARCHIVO if archivo else 'Sin archivo',
                'paginas': detalle.PAGINAS,
                'servicio_id': detalle.SERVICIO_ID,
                'servicio_nombre': servicio_nombre,
                'tamano': detalle.TAMANO,
                'precio': precio,
                'paginas_color': detalle.PAGINAS_COLOR,
                'comentarios': detalle.COMENTARIOS,
                'ruta_descarga': None
            })
            total_paginas += detalle.PAGINAS or 0
            
            # Datos de facturación (si existen)
            if detalle.CANTIDAD and detalle.PRECIO_UNITARIO and detalle.SUBTOTAL:
                detalles_facturacion.append({
                    'cantidad': detalle.CANTIDAD,
                    'precio_unitario': float(detalle.PRECIO_UNITARIO) if detalle.PRECIO_UNITARIO else None,
                    'subtotal': float(detalle.SUBTOTAL)
                })
    # Si no hay detalles, devolvemos arrays vacíos (sin fallback a columnas obsoletas)
    # porque la nueva estructura siempre debe tener detalles.

    return {
        'pedido_id': pedido.ID,
        'estado': pedido.ESTADO,
        'total': float(pedido.TOTAL) if pedido.TOTAL else 0.0,
        'total_paginas': total_paginas,
        'codigo_ticket': pedido.CODIGO_TICKET,
        'referencia_pago': pedido.REFERENCIA_PAGO,
        'comentarios': None,  # ya no existe en pedido; se usa solo para compatibilidad
        'fecha_retiro': pedido.FECHA_RETIRO.strftime('%d/%m/%Y') if pedido.FECHA_RETIRO else None,
        'hora_retiro': pedido.HORA_RETIRO.strftime('%I:%M %p') if pedido.HORA_RETIRO else None,
        'archivos': archivos_info,
        'detalles': detalles_facturacion,
        'usuario': {
            'nombre': usuario.NOMBRE if usuario else None,
            'apellido': usuario.APELLIDO if usuario else None,
            'cedula': usuario.ID_USUARIO if usuario else None,
            'email': usuario.EMAIL if usuario else None,
            'telefono': usuario.TELEFONO if usuario else None
        }
    }

# ------------------------------------------------------------
# PANEL ADMIN – CRONOGRAMA
# ------------------------------------------------------------

@app.route('/admin/cronograma')
@login_required
@admin_required
def admin_cronograma():
    pedidos = Pedido.query.filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso', 'Listo'])
    ).order_by(Pedido.FECHA_RETIRO.asc(), Pedido.HORA_RETIRO.asc()).all()
    
    # Agrupar por fecha
    from collections import defaultdict
    grupos = defaultdict(list)
    for pedido in pedidos:
        if pedido.FECHA_RETIRO:
            usuario = Usuario.query.get(pedido.ID_USUARIO)
            grupos[pedido.FECHA_RETIRO.isoformat()].append({
                'pedido': pedido,
                'usuario': usuario
            })
    
    # Ordenar las fechas
    fechas_ordenadas = sorted(grupos.keys())
    
    # Preparar datos para la plantilla
    cronograma = []
    hoy = date.today()
    for fecha_str in fechas_ordenadas:
        fecha = date.fromisoformat(fecha_str)
        if fecha == hoy:
            etiqueta = "Hoy"
        elif fecha == hoy + timedelta(days=1):
            etiqueta = "Mañana"
        else:
            dias_semana = ['Lunes', 'Martes', 'Miércoles', 'Jueves', 'Viernes', 'Sábado', 'Domingo']
            etiqueta = dias_semana[fecha.weekday()]
        
        cronograma.append({
            'fecha': fecha,
            'etiqueta': etiqueta,
            'fecha_formateada': fecha.strftime('%d/%m/%Y'),
            'pedidos': grupos[fecha_str]
        })
    
    return render_template('admin/cronograma.html', cronograma=cronograma, hoy=hoy)

# ------------------------------------------------------------
# PANEL ADMIN – CONFIGURACION
# ------------------------------------------------------------

@app.route('/admin/configuracion', methods=['GET', 'POST'])
@login_required
@admin_required
def admin_configuracion():
    if request.method == 'POST':
        # Actualizar cada clave con el valor enviado desde el formulario
        for clave, valor in request.form.items():
            if clave.startswith('config_'):
                clave_real = clave[7:]  # quitar prefijo 'config_'
                config = Configuracion.query.filter_by(CLAVE=clave_real).first()
                if config:
                    config.VALOR = valor
        db.session.commit()
        flash('Configuración actualizada.', 'success')
        return redirect(url_for('admin_configuracion'))

    # Cargar toda la configuración actual
    configuraciones = Configuracion.query.all()
    config_dict = {c.CLAVE: c.VALOR for c in configuraciones}
    return render_template('admin/configuracion.html', config=config_dict)

@login_manager.user_loader
def load_user(user_id):
    return db.session.get(Usuario, int(user_id))

def validar_nombre_apellido(texto, campo="Nombre"):
    if not texto:
        return f"El {campo} es obligatorio."
    if len(texto) < 2 or len(texto) > 20:
        return f"El {campo} debe tener entre 2 y 20 caracteres."
    if not re.match(r'^[a-zA-ZáéíóúÁÉÍÓÚñÑ\s]+$', texto):
        return f"El {campo} solo puede contener letras y espacios."
    return None

def validar_cedula(cedula):
    if not cedula:
        return "La cédula es obligatoria."
    cedula = cedula.strip().upper()
    
    # Aceptar solo dígitos de 7 a 12 caracteres
    if re.match(r'^\d{7,12}$', cedula):
        return None
    
    return "Formato de cédula no válido. Use de 7 a 12 dígitos."

def validar_telefono(telefono):
    if not telefono:
        return "El teléfono es obligatorio."
    # Eliminar cualquier carácter no numérico
    telefono = re.sub(r'\D', '', telefono)
    
    # Verificar longitud exacta (11 dígitos para Venezuela)
    if len(telefono) != 11:
        return "El teléfono debe tener exactamente 11 dígitos."
    
    # Lista blanca de prefijos válidos (móviles + un fijo para Caracas)
    prefijos_validos = ['0412', '0422', '0414', '0424', '0416', '0426', '0212']
    
    prefijo = telefono[:4]
    if prefijo not in prefijos_validos:
        return "Prefijo no válido. Solo se permiten números venezolanos (0412, 0414, 0416, 0422, 0424, 0426, 0212)."
    
    return None  # Válido

def validar_email(email):
    if not email:
        return "El correo es obligatorio."
    if not re.match(r'^[^@\s]+@[^@\s]+\.[^@\s]+$', email):
        return "El formato del correo no es válido."
    return None

def validar_contraseña(password):
    errores = []
    if len(password) < 8:
        errores.append("al menos 8 caracteres")
    if not re.search(r"[A-Z]", password):
        errores.append("una mayúscula")
    if not re.search(r"[a-z]", password):
        errores.append("una minúscula")
    if not re.search(r"\d", password):
        errores.append("un número")
    if not re.search(r"[!@#$%^&*(),.?\":{}|<>_\-+=~]", password):
        errores.append("un carácter especial (ej: ! @ # $ %)")

    if errores:
        return f"La contraseña debe contener: {', '.join(errores)}."
    return None

@app.route('/admin/configuracion/vaciar-historial', methods=['POST'])
@login_required
@admin_required
def admin_vaciar_historial():
    try:
        # Eliminar archivos físicos
        carpetas = ['static/uploads/impresion', 'static/uploads/comprobantes']
        for carpeta in carpetas:
            if os.path.exists(carpeta):
                for archivo in os.listdir(carpeta):
                    ruta_archivo = os.path.join(carpeta, archivo)
                    if os.path.isfile(ruta_archivo):
                        os.remove(ruta_archivo)

        # SQL directo para eliminar registros
        db.session.execute(text("DELETE FROM detalle"))
        db.session.execute(text("DELETE FROM archivo_pedido"))
        db.session.execute(text("DELETE FROM pedido"))
        db.session.commit()

        flash('Historial vaciado correctamente.', 'success')
    except Exception as e:
        db.session.rollback()
        flash(f'Error al vaciar el historial: {str(e)}', 'danger')
    
    return redirect(url_for('admin_configuracion'))

@app.route('/admin/limpiar-archivos', methods=['POST'])
@login_required
@admin_required
def admin_limpiar_archivos():
    try:
        eliminados = limpiar_archivos_viejos()
        flash(f'✅ {eliminados} archivos eliminados correctamente.', 'success')
    except Exception as e:
        flash(f'❌ Error: {str(e)}', 'danger')
    return redirect(url_for('admin_configuracion'))

# ------------------------------------------------------------
# PANEL DE OPERADOR
# ------------------------------------------------------------

@app.route('/operador')
@login_required
@operador_required
def panel_operador():
    # Pedidos pendientes de imprimir
    pedidos = Pedido.query.filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso'])
    ).order_by(Pedido.FECHA.asc()).all()

    # Historial reciente (pedidos listos o entregados)
    historial = Pedido.query.filter(
        Pedido.ESTADO.in_(['Listo', 'Entregado'])
    ).order_by(Pedido.FECHA.desc()).limit(20).all()

    # Contador de pedidos no vistos por el operador
    nuevos_operador = Pedido.query.filter_by(VISTO_OPERADOR=False).filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso'])
    ).count()

    # Preparar los datos enriquecidos para las tarjetas
    pedidos_data = []
    for pedido in pedidos:
        usuario = Usuario.query.get(pedido.ID_USUARIO)
        
        # Obtener todos los archivos del pedido
        archivos = ArchivoPedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
        archivos_impresion = [a for a in archivos if 'comprobantes' not in a.RUTA.lower()]
        
        # Obtener detalles (líneas) con su configuración
        detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
        
        # Calcular total de páginas sumando las de cada detalle
        total_paginas = sum(d.PAGINAS for d in detalles) if detalles else (pedido.PAGINAS or 0)
        
        # Obtener lista de servicios (uno por detalle)
        servicios_nombres = []
        for detalle in detalles:
            if detalle.SERVICIO_ID:
                srv = ServicioImpresion.query.get(detalle.SERVICIO_ID)
                if srv:
                    servicios_nombres.append(f"{srv.TITULO} ({detalle.TAMANO})")
        # Si no hay detalles, usar el servicio de la cabecera (fallback)
        if not servicios_nombres and pedido.SERVICIO_ID:
            srv = ServicioImpresion.query.get(pedido.SERVICIO_ID)
            if srv:
                servicios_nombres.append(f"{srv.TITULO} ({pedido.TAMANO})")
        
        pedidos_data.append({
            'pedido': pedido,
            'usuario': usuario,
            'archivos': archivos_impresion,
            'detalles': detalles,
            'total_paginas': total_paginas,
            'servicios': ', '.join(servicios_nombres) if servicios_nombres else 'Sin servicio',
            'referencia_pago': pedido.REFERENCIA_PAGO
        })

    # --- Cronograma (misma lógica que en admin) ---
    from collections import defaultdict
    pedidos_cronograma = Pedido.query.filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso', 'Listo'])
    ).order_by(Pedido.FECHA_RETIRO.asc(), Pedido.HORA_RETIRO.asc()).all()

    grupos = defaultdict(list)
    for pedido in pedidos_cronograma:
        if pedido.FECHA_RETIRO:
            usuario = Usuario.query.get(pedido.ID_USUARIO)
            grupos[pedido.FECHA_RETIRO.isoformat()].append({
                'pedido': pedido,
                'usuario': usuario
            })

    fechas_ordenadas = sorted(grupos.keys())
    cronograma = []
    hoy = date.today()
    for fecha_str in fechas_ordenadas:
        fecha = date.fromisoformat(fecha_str)
        if fecha == hoy:
            etiqueta = "Hoy"
        elif fecha == hoy + timedelta(days=1):
            etiqueta = "Mañana"
        else:
            dias_semana = ['Lunes', 'Martes', 'Miércoles', 'Jueves', 'Viernes', 'Sábado', 'Domingo']
            etiqueta = dias_semana[fecha.weekday()]
        cronograma.append({
            'fecha': fecha,
            'etiqueta': etiqueta,
            'fecha_formateada': fecha.strftime('%d/%m/%Y'),
            'pedidos': grupos[fecha_str]
        })

    # Recuperar mensaje de sesión (si existe)
    mensaje = session.pop('mensaje_operador', None)

    return render_template('operador/dashboard.html',
                           pedidos=pedidos_data,
                           historial=historial,
                           nuevos_operador=nuevos_operador,
                           mensaje=mensaje,
                           cronograma=cronograma)

@app.route('/operador/api/nuevos')
@login_required
@operador_required
def operador_api_nuevos():
    # Total de pedidos en estado "Pago confirmado" o "En proceso"
    total_pendientes = Pedido.query.filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso'])
    ).count()
    
    # Nuevos no vistos por el operador
    nuevos = Pedido.query.filter_by(VISTO_OPERADOR=False).filter(
        Pedido.ESTADO.in_(['Pago confirmado', 'En proceso'])
    ).count()
    
    # Para depuración (opcional)
    print(f"API operador: total={total_pendientes}, nuevos={nuevos}")
    
    return {
        'total_pendientes': total_pendientes,
        'nuevos': nuevos
    }

@app.route('/operador/solicitud/<int:pedido_id>')
@login_required
@operador_required
def operador_detalle(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    usuario = Usuario.query.get(pedido.ID_USUARIO)
    archivos = ArchivoPedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
    detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).all()

    # Solo archivos de impresión
    archivos_impresion = [a for a in archivos if 'comprobantes' not in a.RUTA.lower()]

    # Construir archivos_info desde los detalles
    archivos_info = []
    for detalle in detalles:
        archivo = ArchivoPedido.query.get(detalle.ARCHIVO_ID) if detalle.ARCHIVO_ID else None
        item = {
            'nombre': archivo.NOMBRE_ARCHIVO if archivo else 'Sin archivo',
            'paginas': detalle.PAGINAS,
            'servicio_id': detalle.SERVICIO_ID,
            'tamano': detalle.TAMANO,
            'paginas_color': detalle.PAGINAS_COLOR,
            'comentarios': detalle.COMENTARIOS
        }
        # Enriquecer con nombre de servicio y precio
        if detalle.SERVICIO_ID:
            servicio = ServicioImpresion.query.get(detalle.SERVICIO_ID)
            item['servicio_nombre'] = servicio.TITULO if servicio else 'Desconocido'
            if detalle.TAMANO and servicio:
                tamano_obj = ServicioImpresionTamano.query.filter_by(
                    SERVICIO_ID=detalle.SERVICIO_ID, NOMBRE=detalle.TAMANO
                ).first()
                item['precio'] = float(tamano_obj.PRECIO_BN) if tamano_obj else None
        else:
            item['servicio_nombre'] = None
            item['precio'] = None
        archivos_info.append(item)

    # Calcular total de páginas sumando las de cada detalle
    total_paginas = sum(d.PAGINAS or 0 for d in detalles)

    # Si no hay detalles, no mostrar nada (ya no usamos fallback con columnas obsoletas)
    # Pero si no hay detalles y hay archivos, podemos mostrar un mensaje en la plantilla.

    # Servicio de impresión (para compatibilidad, si se usa en la plantilla)
    servicio = None
    if detalles and detalles[0].SERVICIO_ID:
        servicio = ServicioImpresion.query.get(detalles[0].SERVICIO_ID)

    # Referencia de pago
    referencia_pago = pedido.REFERENCIA_PAGO

    # Marcar como visto
    if not pedido.VISTO_OPERADOR:
        pedido.VISTO_OPERADOR = True
        db.session.commit()

    # Obtener nombre real de cada servicio de impresión
    todos_servicios = {srv.ID: srv.TITULO for srv in ServicioImpresion.query.all()}

    return render_template('operador/detalle.html',
                           pedido=pedido,
                           usuario=usuario,
                           archivos=archivos_impresion,
                           detalles=detalles,
                           servicio=servicio,
                           referencia_pago=referencia_pago,
                           todos_servicios=todos_servicios,
                           archivos_info=archivos_info,
                           total_paginas=total_paginas)

@app.route('/operador/solicitud/<int:pedido_id>/marcar-listo', methods=['POST'])
@login_required
@operador_required
def operador_marcar_listo(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ESTADO in ['Pago confirmado', 'En proceso']:
        # Generar código de ticket si no existe
        if not pedido.CODIGO_TICKET:
            import random, string
            pedido.CODIGO_TICKET = f"TICK-{pedido.ID}-{''.join(random.choices(string.ascii_uppercase + string.digits, k=4))}"

        pedido.ESTADO = 'Listo'
        db.session.commit()

        # Enviar correo al cliente
        usuario = Usuario.query.get(pedido.ID_USUARIO)
        if usuario and usuario.EMAIL:
            enviar_correo_listo(usuario.EMAIL, usuario.NOMBRE, pedido.CODIGO_TICKET, pedido.ID)

        session['mensaje_operador'] = 'Pedido marcado como Listo. Se ha notificado al cliente.'
        return redirect(url_for('panel_operador'))

    return redirect(url_for('panel_operador'))

@app.route('/operador/solicitud/<int:pedido_id>/marcar-entregado', methods=['POST'])
@login_required
@operador_required
def operador_marcar_entregado(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ESTADO in ['Listo', 'Pago confirmado', 'En proceso']:
        pedido.ESTADO = 'Entregado'
        db.session.commit()
        return {'success': True}
    return {'error': 'No se puede marcar como entregado en este estado'}, 400

# ------------------------------------------------------------
# IMPRIMIR
# ------------------------------------------------------------

@app.route('/imprimir', methods=['GET'])
@login_required
def imprimir():
    pedido_id = request.args.get('pedido_id', type=int)
    configurado = request.args.get('configurado') == '1'
    from_config = request.args.get('from_config') == '1'  # ← nuevo
    pedido = None
    archivos_info = []
    cantidad_archivos = 0
    total_paginas = 0

    if pedido_id:
        pedido = Pedido.query.get(pedido_id)
        if pedido and pedido.ID_USUARIO == current_user.ID:
            # Obtener detalles del pedido
            detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
            if detalles:
                for detalle in detalles:
                    archivo = ArchivoPedido.query.get(detalle.ARCHIVO_ID) if detalle.ARCHIVO_ID else None
                    archivos_info.append({
                        'nombre': archivo.NOMBRE_ARCHIVO if archivo else 'Sin archivo',
                        'paginas': detalle.PAGINAS,
                        'servicio_id': detalle.SERVICIO_ID,
                        'tamano': detalle.TAMANO,
                        'paginas_color': detalle.PAGINAS_COLOR,
                        'comentarios': detalle.COMENTARIOS
                    })
                cantidad_archivos = len(detalles)
                total_paginas = sum(d.PAGINAS for d in detalles)
            else:
                # Fallback para pedidos antiguos sin detalles
                archivos = ArchivoPedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
                if archivos:
                    archivos_info = [{'nombre': archivos[0].NOMBRE_ARCHIVO, 'paginas': pedido.PAGINAS or 0}]
                    cantidad_archivos = 1
                    total_paginas = pedido.PAGINAS or 0
        else:
            pedido = None

    return render_template('tienda/imprimir.html',
                           pedido=pedido,
                           archivos_info=archivos_info,
                           cantidad_archivos=cantidad_archivos,
                           total_paginas=total_paginas,
                           configurado=configurado,
                           from_config=from_config,   # ← pasamos al frontend
                           active_page='imprimir')

@app.route('/detectar-paginas', methods=['POST'])
@login_required
def detectar_paginas_ajax():
    if 'archivo' not in request.files:
        return {'error': 'No se recibió archivo'}, 400
    archivo = request.files['archivo']
    if archivo.filename == '':
        return {'error': 'Archivo vacío'}, 400

    filename = secure_filename(archivo.filename)
    ruta_destino = os.path.join(UPLOAD_FOLDER_IMPRESION, filename)
    archivo.save(ruta_destino)

    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    try:
        # 1. Validación MIME e integridad (igual que antes)
        if ext in ('pdf', 'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'pptx'):
            valido, error = validar_mime(ruta_destino, ['pdf', 'imagen', 'pptx'])
        elif ext == 'zip':
            valido, error = validar_mime(ruta_destino, ['zip'])
        else:
            valido, error = False, f"Formato no soportado: .{ext}"
        if not valido:
            os.remove(ruta_destino)
            return {'error': error}, 400

        if ext != 'pdf':
            integro, error_integridad = verificar_integridad(ruta_destino, ext)
            if not integro:
                os.remove(ruta_destino)
                return {'error': error_integridad}, 400

        if ext == 'pdf':
            size = os.path.getsize(ruta_destino)
            if size > 10 * 1024 * 1024:
                os.remove(ruta_destino)
                return {'error': 'El PDF no puede superar los 10 MB.'}, 400
            with open(ruta_destino, 'rb') as f:
                header = f.read(4)
                if not header.startswith(b'%PDF'):
                    os.remove(ruta_destino)
                    return {'error': 'El archivo PDF no es válido (cabecera incorrecta).'}, 400

        # 2. Detección de páginas
        paginas, mensaje = detectar_paginas(ruta_destino, filename)
        if paginas is None:
            os.remove(ruta_destino)
            return {'error': mensaje or 'Formato no soportado'}, 400

        # 3. Crear pedido borrador
        nuevo_pedido = Pedido(
            ID_USUARIO=current_user.ID,
            FECHA=datetime.now(),
            ESTADO='borrador',
            TOTAL=0.0
        )
        db.session.add(nuevo_pedido)
        db.session.flush()  # para obtener ID

        # 4. Crear archivo
        archivo_bd = ArchivoPedido(
            PEDIDO_ID=nuevo_pedido.ID,
            NOMBRE_ARCHIVO=filename,
            RUTA=ruta_destino
        )
        db.session.add(archivo_bd)
        db.session.flush()

        # 5. Crear detalle (línea) para este archivo
        detalle = DetallePedido(
            PEDIDO_ID=nuevo_pedido.ID,
            SERVICIO_ID=None,
            TAMANO=None,
            PAGINAS=paginas,
            PAGINAS_COLOR=None,
            COMENTARIOS=None,
            ARCHIVO_ID=archivo_bd.ID,
            CANTIDAD=paginas,
            PRECIO_UNITARIO=0,
            SUBTOTAL=0
        )
        db.session.add(detalle)
        db.session.commit()

        # 6. Limpiar borradores viejos (eliminando también sus detalles)
        limite = datetime.now() - timedelta(minutes=5)
        pedidos_viejos = Pedido.query.filter(
            Pedido.ID_USUARIO == current_user.ID,
            Pedido.ESTADO == 'borrador',
            Pedido.FECHA < limite
        ).all()
        for p in pedidos_viejos:
            # Eliminar archivos físicos y registros
            for a in ArchivoPedido.query.filter_by(PEDIDO_ID=p.ID).all():
                if os.path.exists(a.RUTA):
                    os.remove(a.RUTA)
                db.session.delete(a)
            # Eliminar detalles
            DetallePedido.query.filter_by(PEDIDO_ID=p.ID).delete()
            # Eliminar pedido
            db.session.delete(p)
        db.session.commit()

        return {
            'success': True,
            'pedido_id': nuevo_pedido.ID,
            'paginas': paginas,
            'mensaje': mensaje or ''
        }

    except Exception as e:
        if os.path.exists(ruta_destino):
            os.remove(ruta_destino)
        db.session.rollback()
        print(f"Error en detectar_paginas_ajax: {e}")
        return {'error': f'Error al procesar el archivo. ({str(e)[:80]})'}, 500

@app.route('/detectar-paginas-multiples', methods=['POST'])
@login_required
def detectar_paginas_multiples():
    import logging
    logging.basicConfig(level=logging.DEBUG)
    logger = logging.getLogger(__name__)

    archivos = request.files.getlist('archivos[]')
    if not archivos or len(archivos) == 0:
        return {'error': 'No se recibieron archivos'}, 400

    docs = 0
    imgs = 0
    for f in archivos:
        ext = f.filename.rsplit('.', 1)[-1].lower() if '.' in f.filename else ''
        if ext in ('pdf', 'pptx', 'zip'):
            docs += 1
        elif ext in ('png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif'):
            imgs += 1
        else:
            return {'error': f'Formato no soportado: {f.filename}'}, 400

    if docs > 10:
        return {'error': 'Máximo 10 documentos permitidos'}, 400
    if imgs > 20:
        return {'error': 'Máximo 20 imágenes permitidas'}, 400

    nuevo_pedido = Pedido(
        ID_USUARIO=current_user.ID,
        FECHA=datetime.now(),
        ESTADO='borrador',
        TOTAL=0.0
    )
    db.session.add(nuevo_pedido)
    db.session.flush()

    total_paginas = 0
    detalles_creados = []

    try:
        for f in archivos:
            filename = secure_filename(f.filename)
            ruta_destino = os.path.join(UPLOAD_FOLDER_IMPRESION, filename)
            logger.debug(f"Guardando archivo: {filename} en {ruta_destino}")
            f.save(ruta_destino)

            ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
            logger.debug(f"Extensión: {ext}")

            # 1. Validación MIME
            if ext in ('pdf', 'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'tif', 'pptx'):
                valido, error = validar_mime(ruta_destino, ['pdf', 'imagen', 'pptx'])
            elif ext == 'zip':
                valido, error = validar_mime(ruta_destino, ['zip'])
            else:
                valido, error = False, f"Formato no soportado: .{ext}"
            if not valido:
                logger.error(f"MIME inválido: {error}")
                os.remove(ruta_destino)
                raise Exception(error)

            # 2. Integridad
            if ext != 'pdf':
                integro, error_integridad = verificar_integridad(ruta_destino, ext)
                if not integro:
                    logger.error(f"Integridad fallida: {error_integridad}")
                    os.remove(ruta_destino)
                    raise Exception(error_integridad)

            # 3. PDF específico
            if ext == 'pdf':
                size = os.path.getsize(ruta_destino)
                if size > 10 * 1024 * 1024:
                    os.remove(ruta_destino)
                    raise Exception('El PDF no puede superar los 10 MB.')
                with open(ruta_destino, 'rb') as fh:
                    header = fh.read(4)
                    if not header.startswith(b'%PDF'):
                        os.remove(ruta_destino)
                        raise Exception('El archivo PDF no es válido (cabecera incorrecta).')

            # 4. Detección de páginas
            paginas, mensaje = detectar_paginas(ruta_destino, filename)
            if paginas is None:
                logger.error(f"Detección de páginas falló para {filename}: {mensaje}")
                os.remove(ruta_destino)
                raise Exception(f'Error en {filename}: {mensaje}')

            total_paginas += paginas
            logger.debug(f"Páginas detectadas: {paginas} para {filename}")

            # 5. Guardar archivo en BD
            archivo_bd = ArchivoPedido(
                PEDIDO_ID=nuevo_pedido.ID,
                NOMBRE_ARCHIVO=filename,
                RUTA=ruta_destino
            )
            db.session.add(archivo_bd)
            db.session.flush()

            # 6. Crear detalle para este archivo
            detalle = DetallePedido(
                PEDIDO_ID=nuevo_pedido.ID,
                SERVICIO_ID=None,
                TAMANO=None,
                PAGINAS=paginas,
                PAGINAS_COLOR=None,
                COMENTARIOS=None,
                ARCHIVO_ID=archivo_bd.ID,
                CANTIDAD=paginas,
                PRECIO_UNITARIO=0,
                SUBTOTAL=0
            )
            db.session.add(detalle)
            detalles_creados.append({
                'nombre': filename,
                'paginas': paginas,
                'servicio_id': None,
                'tamano': None
            })

        db.session.commit()
        logger.info(f"Pedido {nuevo_pedido.ID} creado con {len(detalles_creados)} archivos")

        # Limpiar borradores viejos (con detalles)
        limite = datetime.now() - timedelta(minutes=5)
        pedidos_viejos = Pedido.query.filter(
            Pedido.ID_USUARIO == current_user.ID,
            Pedido.ESTADO == 'borrador',
            Pedido.FECHA < limite
        ).all()
        for p in pedidos_viejos:
            for a in ArchivoPedido.query.filter_by(PEDIDO_ID=p.ID).all():
                if os.path.exists(a.RUTA):
                    os.remove(a.RUTA)
                db.session.delete(a)
            DetallePedido.query.filter_by(PEDIDO_ID=p.ID).delete()
            db.session.delete(p)
        db.session.commit()

        return {
            'success': True,
            'pedido_id': nuevo_pedido.ID,
            'paginas_totales': total_paginas,
            'detalle_archivos': detalles_creados
        }

    except Exception as e:
        # Si ocurrió algún error, hacer rollback y limpiar archivos
        db.session.rollback()
        # Eliminar archivos que pudieron haberse guardado
        if 'ruta_destino' in locals() and os.path.exists(ruta_destino):
            os.remove(ruta_destino)
        # También eliminar archivos de este pedido que ya se hayan guardado en BD
        if 'nuevo_pedido' in locals():
            archivos_guardados = ArchivoPedido.query.filter_by(PEDIDO_ID=nuevo_pedido.ID).all()
            for a in archivos_guardados:
                if os.path.exists(a.RUTA):
                    os.remove(a.RUTA)
                db.session.delete(a)
            DetallePedido.query.filter_by(PEDIDO_ID=nuevo_pedido.ID).delete()
            db.session.delete(nuevo_pedido)
            db.session.commit()
        logger.error(f"Error en detectar_paginas_multiples: {e}", exc_info=True)
        return {'error': f'Error al procesar los archivos. ({str(e)[:80]})'}, 500

@app.route('/cancelar-pedido/<int:pedido_id>', methods=['POST'])
@login_required
def cancelar_pedido(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return {'error': 'No autorizado'}, 403
    if pedido.ESTADO not in ['borrador', 'Pendiente de pago']:
        return {'error': 'No se puede cancelar en este estado'}, 400

    # Eliminar archivos asociados
    archivos = ArchivoPedido.query.filter_by(PEDIDO_ID=pedido_id).all()
    for archivo in archivos:
        if os.path.exists(archivo.RUTA):
            os.remove(archivo.RUTA)
        db.session.delete(archivo)

    # Eliminar detalles
    DetallePedido.query.filter_by(PEDIDO_ID=pedido_id).delete()

    # Eliminar pedido
    db.session.delete(pedido)
    db.session.commit()
    return {'success': True}

@app.route('/configurar/<int:pedido_id>', methods=['GET', 'POST'])
@login_required
def configurar_impresion(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return redirect(url_for('imprimir'))

    if pedido.ESTADO not in ['borrador', 'Pendiente de pago']:
        return redirect(url_for('mis_solicitudes'))

    # Obtener el detalle (asumimos que solo hay uno para pedido único)
    detalle = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).first()
    if not detalle:
        flash('No se encontró el archivo.', 'danger')
        return redirect(url_for('imprimir'))

    # Obtener el archivo asociado
    archivo = ArchivoPedido.query.get(detalle.ARCHIVO_ID) if detalle.ARCHIVO_ID else None
    archivo_nombre = archivo.NOMBRE_ARCHIVO if archivo else 'Sin archivo'
    paginas = detalle.PAGINAS or 1

    servicios = ServicioImpresion.query.filter_by(ACTIVO=True).order_by(ServicioImpresion.TITULO).all()

    # Si se solicita reset, limpiar selección en el detalle
    if request.args.get('reset') == '1':
        detalle.SERVICIO_ID = None
        detalle.TAMANO = None
        detalle.PAGINAS_COLOR = None
        detalle.COMENTARIOS = None
        db.session.commit()
        return redirect(url_for('configurar_impresion', pedido_id=pedido.ID))

    if request.method == 'POST':
        servicio_id = request.form.get('servicio_id', type=int)
        tamano_nombre = request.form.get('tamano_nombre', '').strip()
        comentarios = request.form.get('comentarios', '').strip()
        paginas_color = request.form.get('paginas_color', '').strip() or None

        if not servicio_id or not tamano_nombre:
            flash('Selecciona un servicio y un tamaño.', 'warning')
            return redirect(url_for('configurar_impresion', pedido_id=pedido.ID))

        detalle.SERVICIO_ID = servicio_id
        detalle.TAMANO = tamano_nombre
        detalle.COMENTARIOS = comentarios
        detalle.PAGINAS_COLOR = paginas_color

        db.session.commit()
        return redirect(url_for('programar_retiro', pedido_id=pedido.ID))

    return render_template('tienda/configurar_impresion.html',
                           pedido=pedido,
                           detalle=detalle,
                           servicios=servicios,
                           archivo_nombre=archivo_nombre,
                           paginas=paginas)

@app.route('/configurar-multiple/<int:pedido_id>', methods=['GET', 'POST'])
@login_required
def configurar_impresion_multiple(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return redirect(url_for('imprimir'))

    if pedido.ESTADO not in ['borrador', 'Pendiente de pago']:
        return redirect(url_for('mis_solicitudes'))

    # Obtener detalles
    detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
    if not detalles:
        return redirect(url_for('configurar_impresion', pedido_id=pedido.ID))

    total_paginas = sum(d.PAGINAS for d in detalles)
    servicios = ServicioImpresion.query.filter_by(ACTIVO=True).order_by(ServicioImpresion.TITULO).all()

    # Reset
    if request.args.get('reset') == '1':
        for detalle in detalles:
            detalle.SERVICIO_ID = None
            detalle.TAMANO = None
            detalle.PAGINAS_COLOR = None
            detalle.COMENTARIOS = None
        db.session.commit()
        return redirect(url_for('configurar_impresion_multiple', pedido_id=pedido.ID))

    if request.method == 'POST':
        for i, detalle in enumerate(detalles):
            serv_id = request.form.get(f'servicio_{i}')
            tam = request.form.get(f'tamano_{i}')
            pag_color = request.form.get(f'paginas_color_{i}')
            com = request.form.get(f'comentarios_{i}')
            if serv_id:
                detalle.SERVICIO_ID = int(serv_id)
            if tam:
                detalle.TAMANO = tam
            if pag_color:
                detalle.PAGINAS_COLOR = pag_color
            if com:
                detalle.COMENTARIOS = com
        db.session.commit()
        return redirect(url_for('programar_retiro', pedido_id=pedido.ID))

    return render_template('tienda/configurar_impresion_multiple.html',
                           pedido=pedido,
                           servicios=servicios,
                           detalles=detalles,
                           total_paginas=total_paginas)

@app.route('/api/pedido/<int:pedido_id>')
@login_required
def api_pedido(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return {'error': 'No autorizado'}, 403

    # Obtener detalles del pedido (líneas)
    detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido_id).all()
    archivos_info = []
    total_paginas = 0

    if detalles:
        for detalle in detalles:
            archivo = ArchivoPedido.query.get(detalle.ARCHIVO_ID) if detalle.ARCHIVO_ID else None
            archivos_info.append({
                'nombre': archivo.NOMBRE_ARCHIVO if archivo else 'Sin archivo',
                'paginas': detalle.PAGINAS,
                'servicio_id': detalle.SERVICIO_ID,
                'tamano': detalle.TAMANO,
                'paginas_color': detalle.PAGINAS_COLOR,
                'comentarios': detalle.COMENTARIOS
            })
            total_paginas += detalle.PAGINAS or 0
    else:
        # Fallback para pedidos antiguos sin detalles
        archivos = ArchivoPedido.query.filter_by(PEDIDO_ID=pedido_id).all()
        if archivos:
            archivos_info = [{
                'nombre': archivos[0].NOMBRE_ARCHIVO,
                'paginas': pedido.PAGINAS or 0
            }]
            total_paginas = pedido.PAGINAS or 0

    # Obtener servicio si existe (para compatibilidad)
    servicio_nombre = None
    if pedido.SERVICIO_ID:
        srv = ServicioImpresion.query.get(pedido.SERVICIO_ID)
        if srv:
            servicio_nombre = srv.TITULO

    return {
        'pedido_id': pedido.ID,
        'estado': pedido.ESTADO,
        'total': float(pedido.TOTAL),
        'archivos': archivos_info,          # ← lista completa de archivos
        'total_paginas': total_paginas,     # ← suma de todas las páginas
        'fecha_retiro': pedido.FECHA_RETIRO.strftime('%Y-%m-%d') if pedido.FECHA_RETIRO else '',
        'hora_retiro': pedido.HORA_RETIRO.strftime('%H:%M') if pedido.HORA_RETIRO else '',
        'servicio': servicio_nombre,
        'tamano': pedido.TAMANO,
        'archivo_nombre': archivos_info[0]['nombre'] if archivos_info else 'Sin archivo',
        'paginas': total_paginas,           # ← para compatibilidad con el frontend
        'detalle_archivos': None            # ya no usamos JSON
    }

@app.route('/programar-retiro/<int:pedido_id>', methods=['GET', 'POST'])
@login_required
def programar_retiro(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return redirect(url_for('imprimir'))

    if pedido.ESTADO not in ['borrador', 'Pendiente de pago']:
        return redirect(url_for('mis_solicitudes'))

    # Obtener detalles
    detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
    if not detalles:
        flash('No hay archivos configurados para este pedido.', 'danger')
        return redirect(url_for('imprimir'))

    # Calcular total de páginas
    total_paginas = sum(d.PAGINAS for d in detalles)

    # Preparar información de archivos para la plantilla
    archivos_info = []
    for detalle in detalles:
        archivo = ArchivoPedido.query.get(detalle.ARCHIVO_ID) if detalle.ARCHIVO_ID else None
        archivos_info.append({
            'nombre': archivo.NOMBRE_ARCHIVO if archivo else 'Sin archivo',
            'paginas': detalle.PAGINAS,
            'servicio_id': detalle.SERVICIO_ID,
            'tamano': detalle.TAMANO,
            'paginas_color': detalle.PAGINAS_COLOR,
            'comentarios': detalle.COMENTARIOS
        })

    if request.method == 'POST':
        fecha_retiro = request.form.get('fecha_retiro')
        hora_retiro = request.form.get('hora_retiro')

        if not fecha_retiro or not hora_retiro:
            flash('Selecciona fecha y hora.', 'warning')
            return redirect(url_for('programar_retiro', pedido_id=pedido.ID))

        try:
            fecha_elegida = datetime.strptime(fecha_retiro, '%Y-%m-%d').date()
            hora_elegida = datetime.strptime(hora_retiro, '%H:%M').time()
        except ValueError:
            flash('Formato de fecha u hora no válido.', 'danger')
            return redirect(url_for('programar_retiro', pedido_id=pedido.ID))

        # Validaciones de horario (igual que antes)
        if fecha_elegida.weekday() == 6:
            flash('No se pueden programar retiros los domingos.', 'warning')
            return redirect(url_for('programar_retiro', pedido_id=pedido.ID))

        if fecha_elegida.weekday() == 5:
            horarios_validos = [(9,0),(9,30),(10,0),(10,30),(11,0),(11,30),(12,0)]
            permitido = any(h == hora_elegida.hour and m == hora_elegida.minute for h, m in horarios_validos)
            if not permitido:
                flash('Los sábados solo se puede retirar de 9:00 AM a 12:00 PM.', 'warning')
                return redirect(url_for('programar_retiro', pedido_id=pedido.ID))
        else:
            manana = [(7,0),(7,30),(8,0),(8,30),(9,0),(9,30),(10,0),(10,30),(11,0),(11,30),(12,0)]
            tarde  = [(13,0),(13,30),(14,0),(14,30),(15,0),(15,30),(16,0),(16,30),(17,0),(17,30)]
            permitido = any((h == hora_elegida.hour and m == hora_elegida.minute) for h, m in manana + tarde)
            if not permitido:
                flash('Horario no válido. L-V: 7:00 AM - 12:00 PM y 1:00 PM - 5:30 PM.', 'warning')
                return redirect(url_for('programar_retiro', pedido_id=pedido.ID))

        fecha_minima = datetime.now().date() + timedelta(days=1)
        fecha_limite = datetime.now().date() + timedelta(days=7)
        if fecha_elegida < fecha_minima or fecha_elegida > fecha_limite:
            flash('La fecha de retiro debe ser entre mañana y los próximos 7 días.', 'warning')
            return redirect(url_for('programar_retiro', pedido_id=pedido.ID))

        # ======= CALCULAR TOTAL DESDE DETALLES =======
        total = 0
        for detalle in detalles:
            if not detalle.SERVICIO_ID or not detalle.TAMANO:
                flash('Todos los archivos deben tener un servicio y tamaño seleccionados.', 'danger')
                return redirect(url_for('configurar_impresion_multiple', pedido_id=pedido.ID))

            servicio = ServicioImpresion.query.get(detalle.SERVICIO_ID)
            if not servicio:
                continue

            tamano_obj = ServicioImpresionTamano.query.filter_by(
                SERVICIO_ID=detalle.SERVICIO_ID, NOMBRE=detalle.TAMANO
            ).first()
            if not tamano_obj:
                continue

            paginas = detalle.PAGINAS or 0
            precio_bn = float(tamano_obj.PRECIO_BN)
            precio_color = float(tamano_obj.PRECIO_COLOR) if servicio.ES_MIXTO else precio_bn
            subtotal = 0

            if servicio.ES_MIXTO and detalle.PAGINAS_COLOR:
                try:
                    paginas_color = int(detalle.PAGINAS_COLOR)
                except ValueError:
                    paginas_color = 0
                if paginas_color > paginas:
                    paginas_color = paginas
                paginas_bn = paginas - paginas_color
                subtotal = (precio_bn * paginas_bn) + (precio_color * paginas_color)
            else:
                subtotal = precio_bn * paginas

            detalle.PRECIO_UNITARIO = precio_bn if not servicio.ES_MIXTO else precio_color
            detalle.SUBTOTAL = subtotal
            total += subtotal

        total = round(total, 2)

        pedido.FECHA_RETIRO = fecha_elegida
        pedido.HORA_RETIRO = hora_elegida
        pedido.TOTAL = total
        pedido.ESTADO = 'Pendiente de pago'

        db.session.commit()
        return redirect(url_for('pagar_impresion', pedido_id=pedido.ID))

    # GET
    return render_template('tienda/programar_retiro.html',
                           pedido=pedido,
                           detalles=detalles,
                           archivos_info=archivos_info,
                           paginas=total_paginas)

@app.route('/pagar-impresion/<int:pedido_id>', methods=['GET', 'POST'])
@login_required
def pagar_impresion(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if request.method == 'POST':
        referencia = request.form.get('referencia')
        comprobante = request.files.get('comprobante')
        if comprobante:
            filename = secure_filename(comprobante.filename)
            filepath = os.path.join(UPLOAD_FOLDER_COMPROBANTES, filename)  # antes UPLOAD_FOLDER_IMPRESION
            comprobante.save(filepath)
            archivo = ArchivoPedido(
                PEDIDO_ID=pedido.ID,
                NOMBRE_ARCHIVO=filename,
                RUTA=filepath
            )
            db.session.add(archivo)
        pedido.REFERENCIA_PAGO = referencia
        pedido.ESTADO = 'Esperando validación'
        db.session.commit()
        return redirect(url_for('espera_validacion', pedido_id=pedido.ID))

    # Obtener datos de PagoMóvil desde la tabla configuracion
    config_pago = {}
    configs = Configuracion.query.filter(Configuracion.CLAVE.startswith('pago_movil')).all()
    for c in configs:
        config_pago[c.CLAVE] = c.VALOR

    return render_template('tienda/pagar_impresion.html', pedido=pedido, config_pago=config_pago)

@app.route('/cancelar-impresion/<int:pedido_id>', methods=['POST'])
@login_required
def cancelar_impresion(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return {'error': 'No autorizado'}, 403

    if pedido.ESTADO in ['Cancelado', 'Pago confirmado', 'Listo', 'Entregado']:
        return {'error': 'El pedido ya no puede ser cancelado.'}, 400

    pedido.ESTADO = 'Cancelado'
    db.session.commit()

    return {'ok': True}

@app.route('/espera-validacion/<int:pedido_id>')
@login_required
def espera_validacion(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return redirect(url_for('imprimir'))
    return render_template('tienda/espera_validacion.html', pedido=pedido)

@app.route('/api/estado-pedido/<int:pedido_id>')
@login_required
def api_estado_pedido(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return {'error': 'No autorizado'}, 403
    return {'estado': pedido.ESTADO}

@app.route('/ticket/<int:pedido_id>')
@login_required
def ticket_impresion(pedido_id):
    pedido = Pedido.query.get_or_404(pedido_id)
    if pedido.ID_USUARIO != current_user.ID:
        return redirect(url_for('imprimir'))

    if pedido.ESTADO not in ['Pago confirmado', 'Listo', 'Entregado']:
        return redirect(url_for('espera_validacion', pedido_id=pedido.ID))

    usuario = Usuario.query.get(pedido.ID_USUARIO)
    if not pedido.CODIGO_TICKET:
        # Generar por si acaso (redundante)
        import random, string
        codigo = f"TICK-{pedido.ID}-{''.join(random.choices(string.ascii_uppercase + string.digits, k=4))}"
        pedido.CODIGO_TICKET = codigo
        db.session.commit()
    return render_template('tienda/ticket_impresion.html', pedido=pedido, usuario=usuario)

# ------------------------------------------------------------
# SOLICITUDES
# ------------------------------------------------------------

@app.route('/mis-solicitudes')
@login_required
def mis_solicitudes():
    pedidos = Pedido.query.filter_by(ID_USUARIO=current_user.ID).order_by(Pedido.FECHA.desc()).all()
    
    pedidos_data = []
    for pedido in pedidos:
        archivos = ArchivoPedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
        detalles = DetallePedido.query.filter_by(PEDIDO_ID=pedido.ID).all()
        
        # Calcular total de páginas sumando las de cada detalle
        total_paginas = sum(d.PAGINAS or 0 for d in detalles)
        
        # Construir información de servicios desde los detalles
        servicio_nombre = None
        tamano = None
        paginas_color = None
        comentarios_generales = None
        
        if detalles:
            servicios_nombres = []
            tamanos = []
            colores = []
            comentarios_list = []
            for detalle in detalles:
                if detalle.SERVICIO_ID:
                    srv = ServicioImpresion.query.get(detalle.SERVICIO_ID)
                    if srv:
                        servicios_nombres.append(f"{srv.TITULO} · {detalle.TAMANO or 'Sin tamaño'}")
                if detalle.TAMANO:
                    tamanos.append(detalle.TAMANO)
                if detalle.PAGINAS_COLOR:
                    colores.append(detalle.PAGINAS_COLOR)
                if detalle.COMENTARIOS:
                    comentarios_list.append(detalle.COMENTARIOS)
            servicio_nombre = ', '.join(servicios_nombres) if servicios_nombres else None
            tamano = ', '.join(tamanos) if tamanos else None
            paginas_color = ', '.join(colores) if colores else None
            comentarios_generales = ', '.join(comentarios_list) if comentarios_list else None
        else:
            # Fallback para pedidos antiguos sin detalles (por si acaso)
            if pedido.SERVICIO_ID:
                srv = ServicioImpresion.query.get(pedido.SERVICIO_ID)
                if srv:
                    servicio_nombre = srv.TITULO
            tamano = pedido.TAMANO
            paginas_color = pedido.PAGINAS_COLOR
            comentarios_generales = pedido.COMENTARIOS
            if not total_paginas:
                total_paginas = pedido.PAGINAS or 0
        
        pedidos_data.append({
            'pedido': pedido,
            'archivos': archivos,
            'detalles': detalles,
            'servicio_nombre': servicio_nombre,
            'tamano': tamano,
            'paginas_color': paginas_color,
            'comentarios': comentarios_generales,
            'total_paginas': total_paginas
        })
    
    return render_template('tienda/mis_solicitudes.html', pedidos=pedidos_data, active_page='solicitudes')

# ------------------------------------------------------------
# CERRAR SESION
# ------------------------------------------------------------

@app.route('/logout')
@login_required
def logout():
    logout_user()
    # Limpiar todos los mensajes flash anteriores para que no aparezcan en el login
    session.clear()
    return redirect(url_for('login', logout=1))

# ------------------------------------------------------------
# TERMINOS
# ------------------------------------------------------------

@app.route('/politicas')
def politicas():
    config = Configuracion.query.filter_by(CLAVE='politicas').first()
    texto = config.VALOR if config else 'Políticas y términos no configurados aún.'
    return render_template('politicas.html', texto=texto)

# ------------------------------------------------------------
# MI PERFIL
# ------------------------------------------------------------

@app.route('/mi-perfil', methods=['GET', 'POST'])
@login_required
def mi_perfil():
    user = current_user

    if request.method == 'POST':
        email = request.form.get('email', '').strip().lower()
        telefono = request.form.get('telefono', '').strip()
        pregunta1 = request.form.get('pregunta1', '').strip()
        respuesta1 = request.form.get('respuesta1', '').strip()
        pregunta2 = request.form.get('pregunta2', '').strip()
        respuesta2 = request.form.get('respuesta2', '').strip()

        field_errors = {}

        # Validar email
        if not email or '@' not in email or '.' not in email.split('@')[-1]:
            field_errors['email'] = 'Ingresa un correo válido.'
        else:
            # Verificar que no esté en uso por otro usuario
            existente = Usuario.query.filter(Usuario.EMAIL == email, Usuario.ID != user.ID).first()
            if existente:
                field_errors['email'] = 'Este correo ya está en uso por otro usuario.'

        # Validar teléfono
        err_tel = validar_telefono(telefono)
        if err_tel:
            field_errors['telefono'] = err_tel
        else:
            existente = Usuario.query.filter(Usuario.TELEFONO == telefono, Usuario.ID != user.ID).first()
            if existente:
                field_errors['telefono'] = 'Este número de teléfono ya está registrado.'
                
        # Validar preguntas de seguridad
        if pregunta1:
            pregunta1 = pregunta1.strip()
            if len(pregunta1) < 4:
                field_errors['pregunta1'] = 'La pregunta debe tener al menos 4 caracteres.'
            elif len(pregunta1) > 20:
                field_errors['pregunta1'] = 'La pregunta no puede tener más de 20 caracteres.'
            elif not re.match(r'^[a-zA-ZáéíóúÁÉÍÓÚñÑ\s?]+$', pregunta1):
                field_errors['pregunta1'] = 'Solo se permiten letras y espacios.'
        if respuesta1:
            respuesta1 = respuesta1.strip()
            if len(respuesta1) > 20:
                field_errors['respuesta1'] = 'La respuesta no puede tener más de 20 caracteres.'

        if pregunta2:
            pregunta2 = pregunta2.strip()
            if len(pregunta2) < 4:
                field_errors['pregunta2'] = 'La pregunta debe tener al menos 4 caracteres.'
            elif len(pregunta2) > 20:
                field_errors['pregunta2'] = 'La pregunta no puede tener más de 20 caracteres.'
            elif not re.match(r'^[a-zA-ZáéíóúÁÉÍÓÚñÑ\s?]+$', pregunta2):
                field_errors['pregunta2'] = 'Solo se permiten letras y espacios.'
        if respuesta2:
            respuesta2 = respuesta2.strip()
            if len(respuesta2) > 20:
                field_errors['respuesta2'] = 'La respuesta no puede tener más de 20 caracteres.'

        # Validar que las preguntas no sean iguales entre sí
        if pregunta1 and pregunta2 and pregunta1.lower() == pregunta2.lower():
            field_errors['pregunta2'] = 'Las preguntas no pueden ser iguales.'

        # Validar que las respuestas no sean iguales entre sí
        if respuesta1 and respuesta2 and respuesta1.lower() == respuesta2.lower():
            field_errors['respuesta2'] = 'Las respuestas no pueden ser iguales.'

        # Validar que las preguntas no sean iguales
        p1 = request.form.get('pregunta1', '').strip()
        p2 = request.form.get('pregunta2', '').strip()
        if p1 and p2 and p1.lower() == p2.lower():
            field_errors['pregunta2'] = 'Las preguntas de seguridad no pueden ser iguales.'

        # Validar que las respuestas no sean iguales
        r1 = request.form.get('respuesta1', '').strip()
        r2 = request.form.get('respuesta2', '').strip()
        if r1 and r2 and r1.lower() == r2.lower():
            field_errors['respuesta2'] = 'Las respuestas de seguridad no pueden ser iguales.'

        if field_errors:
            return render_template('mi_perfil.html', form_data=request.form, field_errors=field_errors)

        # Actualizar datos del usuario
        user.EMAIL = email
        user.TELEFONO = telefono
        if pregunta1:
            user.PREGUNTA1 = pregunta1
        if respuesta1:
            user.RESPUESTA1 = respuesta1.strip().upper()
        if pregunta2:
            user.PREGUNTA2 = pregunta2
        if respuesta2:
            user.RESPUESTA2 = respuesta2.strip().upper()

        db.session.commit()
        flash('Perfil actualizado correctamente.', 'success')
        return redirect(url_for('mi_perfil'))

    # GET: mostrar formulario con datos actuales
    return render_template('mi_perfil.html', form_data=None, field_errors={})

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)